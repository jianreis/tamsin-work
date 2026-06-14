#!/usr/bin/env python3
"""
build-northstar-deck.py
=======================
Re-runnable, offline, idempotent generator for the consulting-grade executive deck

    deliverables/06-future-state-northstar-deck.pptx

"Future-State Scripting — A North Star for Dis-Chem"  (Bigly Labs x Dis-Chem)

Narrative arc (Deloitte/Bain style): problem -> rigorous method -> evidence ->
principles & people it is grounded in -> the North Star experience -> what makes
it real -> what's next. Every analytical claim is cited to the ground-truth vault.

Dependencies: python-pptx 1.0.2 (already installed). No network, no external
fonts/images. Embeds ONLY the four staged PNGs in deck-assets/.

Run:  python3 deliverables/_shared/build-northstar-deck.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Paths
# ----------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "deck-assets")
OUT = os.path.abspath(os.path.join(HERE, "..", "06-future-state-northstar-deck.pptx"))

IMG_BOARD = os.path.join(ASSETS, "journey-board.png")
IMG_MOMENT = os.path.join(ASSETS, "journey-moment-expanded.png")
IMG_CASH = os.path.join(ASSETS, "journey-cash.png")
IMG_CABINET = os.path.join(ASSETS, "medicine-cabinet.png")

# ----------------------------------------------------------------------------
# Palette (Dis-Chem)
# ----------------------------------------------------------------------------
GREEN = RGBColor(0x1A, 0x9B, 0x4A)   # primary accent / dividers
GREEN_DK = RGBColor(0x12, 0x6E, 0x35)
LIME = RGBColor(0xB5, 0xD3, 0x34)    # sparing accent
INK = RGBColor(0x1D, 0x1D, 0x1D)     # headings / dark text
SLATE = RGBColor(0x52, 0x52, 0x5B)   # body text
MINT = RGBColor(0xE8, 0xF5, 0xEC)    # light fill
MINT_MD = RGBColor(0xBF, 0xE3, 0xCF)  # mid mint fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xFA, 0xFB, 0xFA)   # near-white bg
LINE = RGBColor(0xD7, 0xDD, 0xD9)    # hairline
AMBER = RGBColor(0xB0, 0x6A, 0x00)   # reform / caution accent
CARD_BG = RGBColor(0xF4, 0xF8, 0xF5)

FONT = "Calibri"
FONT_LT = "Calibri Light"

# Slide geometry (16:9)
SW = Inches(13.333)
SH = Inches(7.5)

FOOTER = ("Concept for review — not an official Dis-Chem artefact  ·  "
          "Bigly Labs × Dis-Chem")

# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------

def _set_bg(slide, color):
    """Solid slide background via a full-bleed rectangle sent to back."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    # send to back
    sp = rect._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return rect


def _no_autofit(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0):
    """runs: list of paragraphs; each paragraph is a list of (text, opts) tuples.
    opts: dict with size, bold, italic, color, font, underline."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        for text, opts in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", FONT)
            f.size = Pt(opts.get("size", 14))
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
            f.underline = opts.get("underline", False)
            f.color.rgb = opts.get("color", INK)
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
             shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    return sp


def _soft_shadow(shape):
    """Add a subtle outer shadow to a shape."""
    spPr = shape._element.spPr
    # remove existing effectLst
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    shdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': '60000', 'dist': '25000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = shdw.makeelement(qn('a:srgbClr'), {'val': '1D1D1D'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '22000'})
    clr.append(alpha)
    shdw.append(clr)
    effectLst.append(shdw)
    spPr.append(effectLst)


def shape_text(shape, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
               margins=0.06):
    tf = shape.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margins)
    tf.margin_right = Inches(margins)
    tf.margin_top = Inches(margins * 0.6)
    tf.margin_bottom = Inches(margins * 0.6)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.0
        for text, opts in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", FONT)
            f.size = Pt(opts.get("size", 12))
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
            f.color.rgb = opts.get("color", INK)
    return shape


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ----------------------------------------------------------------------------
# Slide chrome (content slides)
# ----------------------------------------------------------------------------
PAGE = {"n": 0}


def content_chrome(slide, kicker, title, takeaway, cite=None):
    """Standard content-slide header: kicker, title, accent rule, so-what strip."""
    _set_bg(slide, WHITE)
    # top accent bar
    add_rect(slide, 0, 0, SW, Inches(0.12), fill=GREEN)
    # kicker
    add_text(slide, Inches(0.6), Inches(0.34), Inches(11.5), Inches(0.3),
             [[(kicker.upper(), {"size": 11, "bold": True, "color": GREEN,
                                 "font": FONT})]])
    # title
    add_text(slide, Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.9),
             [[(title, {"size": 27, "bold": True, "color": INK})]],
             line_spacing=0.98)
    # so-what strip (takeaway) — mint band with green left tab
    sy = Inches(1.62)
    add_rect(slide, Inches(0.6), sy, Inches(0.09), Inches(0.52), fill=GREEN)
    band = add_rect(slide, Inches(0.69), sy, Inches(12.04), Inches(0.52), fill=MINT)
    shape_text(band, [[("So what:  ", {"size": 12.5, "bold": True, "color": GREEN_DK}),
                       (takeaway, {"size": 12.5, "bold": False, "color": INK})]],
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margins=0.16)
    if cite:
        add_text(slide, Inches(0.6), Inches(7.06), Inches(8.0), Inches(0.3),
                 [[("Source: ", {"size": 8, "italic": True, "color": SLATE}),
                   (cite, {"size": 8, "italic": True, "color": GREEN_DK})]])
    _footer(slide)


def _footer(slide):
    PAGE["n"] += 1
    add_rect(slide, Inches(0.0), Inches(7.18), SW, Pt(0.75), fill=LINE)
    add_text(slide, Inches(0.6), Inches(7.24), Inches(11.0), Inches(0.24),
             [[(FOOTER, {"size": 7.5, "color": SLATE})]])
    add_text(slide, Inches(12.4), Inches(7.24), Inches(0.7), Inches(0.24),
             [[(str(PAGE["n"]), {"size": 8, "bold": True, "color": GREEN_DK})]],
             align=PP_ALIGN.RIGHT)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ----------------------------------------------------------------------------
# Section divider
# ----------------------------------------------------------------------------

def section_divider(prs, num, title, subtitle):
    slide = new_slide(prs)
    _set_bg(slide, GREEN)
    # large faint "+" motif (Dis-Chem mark) top-right
    plus_v = add_rect(slide, Inches(11.7), Inches(0.7), Inches(0.5), Inches(2.4),
                      fill=GREEN_DK)
    plus_h = add_rect(slide, Inches(10.75), Inches(1.65), Inches(2.4), Inches(0.5),
                      fill=GREEN_DK)
    add_text(slide, Inches(0.9), Inches(2.7), Inches(2.0), Inches(0.6),
             [[(num, {"size": 40, "bold": True, "color": LIME,
                      "font": FONT_LT})]])
    add_rect(slide, Inches(0.95), Inches(3.55), Inches(0.9), Inches(0.06), fill=LIME)
    add_text(slide, Inches(0.9), Inches(3.75), Inches(11.0), Inches(1.4),
             [[(title, {"size": 34, "bold": True, "color": WHITE})]],
             line_spacing=1.0)
    add_text(slide, Inches(0.95), Inches(5.0), Inches(10.6), Inches(1.0),
             [[(subtitle, {"size": 14, "color": MINT_MD})]], line_spacing=1.15)
    PAGE["n"] += 1
    add_text(slide, Inches(12.4), Inches(7.24), Inches(0.7), Inches(0.24),
             [[(str(PAGE["n"]), {"size": 8, "bold": True, "color": MINT_MD})]],
             align=PP_ALIGN.RIGHT)
    return slide


# ============================================================================
# SLIDE BUILDERS
# ============================================================================

def slide_title(prs):
    slide = new_slide(prs)
    _set_bg(slide, WHITE)
    # left green panel
    add_rect(slide, 0, 0, Inches(0.45), SH, fill=GREEN)
    add_rect(slide, Inches(0.45), 0, Inches(0.08), SH, fill=LIME)
    # big "+" watermark
    add_rect(slide, Inches(10.55), Inches(0.9), Inches(0.62), Inches(2.9),
             fill=MINT_MD)
    add_rect(slide, Inches(9.4), Inches(2.04), Inches(2.9), Inches(0.62),
             fill=MINT_MD)
    add_text(slide, Inches(1.0), Inches(1.5), Inches(9.0), Inches(0.4),
             [[("BIGLY LABS  ×  DIS-CHEM", {"size": 13, "bold": True,
                                                 "color": GREEN})]])
    add_text(slide, Inches(0.95), Inches(2.25), Inches(10.6), Inches(2.2),
             [[("Future-State Scripting", {"size": 50, "bold": True,
                                            "color": INK})],
              [("A North Star for Dis-Chem", {"size": 50, "bold": True,
                                              "color": GREEN})]],
             line_spacing=1.0)
    add_text(slide, Inches(1.0), Inches(4.55), Inches(10.8), Inches(0.9),
             [[("Defining the prescription experience we want to build — an ",
                {"size": 16, "color": SLATE}),
               ("experience definition, not yet a blueprint.",
                {"size": 16, "italic": True, "color": INK})]],
             line_spacing=1.15)
    # concept note ribbon
    rb = add_rect(slide, Inches(1.0), Inches(5.5), Inches(8.4), Inches(0.5),
                  fill=MINT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(rb, [[("⚠  Concept for review — not an official Dis-Chem artefact.  "
                      "Generated 2026-06-14 from the Bigly ground-truth vault.",
                      {"size": 10.5, "bold": True, "color": GREEN_DK})]],
               anchor=MSO_ANCHOR.MIDDLE, margins=0.14)
    add_text(slide, Inches(1.0), Inches(6.7), Inches(10.0), Inches(0.4),
             [[("Executive deck for Dis-Chem & Bigly leadership  ·  June 2026",
                {"size": 11, "color": SLATE})]])
    notes(slide, "Welcome. This deck takes you through how we built a defensible "
                 "North Star for prescription ('scripting') at Dis-Chem — from the "
                 "problem, through a rigorous evidence-backed method, to the future-"
                 "state experience itself. Frame it up front: this is an experience "
                 "definition we want leadership to agree on, NOT yet a service "
                 "blueprint — that is the next step. Everything here is cited to our "
                 "ground-truth research vault.")
    PAGE["n"] += 1
    return slide


def slide_exec_summary(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Executive summary", "The North Star in one slide",
        "Price is regulated and fixed — so the durable advantage is the experience "
        "around the medicine: a calm, trusted, self-running scripting service for "
        "both South Africas.",
        cite="deliverables/05-tobe-future-state-journey.html · pillars.md")
    # Where we are taking scripting (left) | Why it matters (right)
    top = Inches(2.45)
    colw = Inches(5.95)
    gap = Inches(0.25)
    lx = Inches(0.6)
    rx = lx + colw + gap
    # Left card
    add_rect(slide, lx, top, colw, Inches(4.0), fill=CARD_BG, line=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_text(slide, lx + Inches(0.3), top + Inches(0.22), colw - Inches(0.6),
             Inches(0.4),
             [[("WHERE WE'RE TAKING SCRIPTING", {"size": 11.5, "bold": True,
                                                 "color": GREEN})]])
    bullets_l = [
        "The household medicine manager — one account, every family member, repeats that simply run.",
        "From a transactional, paper-and-queue errand to a managed care relationship.",
        "One service spine, two first-class funded paths — insured and cash.",
        "Automation does the toil; a named pharmacist stays visible and in control.",
    ]
    _bullet_block(slide, lx + Inches(0.3), top + Inches(0.62), colw - Inches(0.6),
                  bullets_l, size=12.5)
    # Right card (green)
    add_rect(slide, rx, top, colw, Inches(4.0), fill=GREEN,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_text(slide, rx + Inches(0.3), top + Inches(0.22), colw - Inches(0.6),
             Inches(0.4),
             [[("WHY IT MATTERS", {"size": 11.5, "bold": True, "color": LIME})]])
    bullets_r = [
        ("65% negative", " pharmacy sentiment — Dis-Chem wins on retail price but is losing on pharmacy trust."),
        ("Customers defect", " to local pharmacies, willing to pay more for relationship care."),
        ("Chronic repeats", " are the recurring relationship that compounds loyalty — the engine to own."),
        ("Buy-in first", " — agree the experience here, then blueprint and build it."),
    ]
    yy = top + Inches(0.7)
    for lead, rest in bullets_r:
        add_text(slide, rx + Inches(0.3), yy, colw - Inches(0.6), Inches(0.8),
                 [[("▸ ", {"size": 12.5, "bold": True, "color": LIME}),
                   (lead, {"size": 12.5, "bold": True, "color": WHITE}),
                   (rest, {"size": 12.5, "color": MINT_MD})]],
                 line_spacing=1.05)
        yy += Inches(0.86)
    notes(slide, "The whole argument on one page. Where we are taking scripting: a "
                 "self-running, household-level medicine service that turns a "
                 "transactional errand into a managed care relationship — built for "
                 "both the insured 16% and the cash 84% on one spine. Why it matters: "
                 "pharmacy sentiment runs 65% negative even as the retail brand wins "
                 "on price, and customers are defecting to local pharmacies. Price is "
                 "regulated, so experience is the only durable differentiator. Ask of "
                 "this deck: agree the experience, then we blueprint and build.")
    return slide


def _bullet_block(slide, x, y, w, items, size=12.5, color=INK, gap=0.82,
                  marker="▸ ", marker_color=GREEN):
    yy = y
    for it in items:
        add_text(slide, x, yy, w, Inches(0.78),
                 [[(marker, {"size": size, "bold": True, "color": marker_color}),
                   (it, {"size": size, "color": color})]], line_spacing=1.05)
        yy += Inches(gap)


def slide_why_now(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Why now — the problem", "Pharmacy pain is a trust problem, not an availability problem",
        "The medicine is usually there; what fails is the experience around it — "
        "and in healthcare, trust failures don't get a second chance.",
        cite="wiki/dischem/dischem-social-sentiment-trust-journey · wiki/dischem/pharmacy-complaints-pain-points")
    top = Inches(2.4)
    # Left: trust journey arc (Calm -> Defection)
    add_text(slide, Inches(0.6), top, Inches(7.3), Inches(0.3),
             [[("THE PATIENT TRUST JOURNEY", {"size": 11, "bold": True, "color": GREEN}),
               ("   (social listening, 1,270 coded comments)", {"size": 10, "italic": True, "color": SLATE})]])
    stages = [
        ("Calm", "Assumes it will work", MINT_MD),
        ("Trigger", "Re-keying, missed windows", RGBColor(0xCB, 0xD8, 0x6E)),
        ("Acceleration", "Errors feel life-threatening", LIME),
        ("Agitation", "Can't reach anyone", RGBColor(0xE0, 0xA9, 0x3A)),
        ("Defection", "Moves scripts away", RGBColor(0xC9, 0x6A, 0x2E)),
    ]
    ax = Inches(0.6)
    aw = Inches(1.42)
    ah = Inches(1.55)
    ay = top + Inches(0.42)
    for i, (st, sub, col) in enumerate(stages):
        chev = add_rect(slide, ax + i * (aw - Inches(0.06)), ay, aw, Inches(0.92),
                        fill=col, shape=MSO_SHAPE.CHEVRON)
        dark = col in (LIME, RGBColor(0xCB, 0xD8, 0x6E), MINT_MD)
        shape_text(chev, [[(st, {"size": 11.5, "bold": True,
                                 "color": INK if dark else WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.04)
        add_text(slide, ax + i * (aw - Inches(0.06)), ay + Inches(1.0),
                 aw - Inches(0.1), Inches(0.7),
                 [[(sub, {"size": 8.5, "color": SLATE})]],
                 align=PP_ALIGN.CENTER, line_spacing=0.95)
    # quote under arc
    qb = add_rect(slide, Inches(0.6), ay + ah + Inches(0.55), Inches(7.3), Inches(1.55),
                  fill=CARD_BG, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(qb, [
        [("“I have been pulling at my hair trying to order my medication "
          "via the App or even using the WhatsApp service.”",
          {"size": 12.5, "italic": True, "color": INK})],
        [("", {"size": 5})],
        [("“Processed on my medical aid and out for delivery. I did not "
          "receive it.”   ", {"size": 11.5, "italic": True, "color": SLATE}),
         ("— real Dis-Chem customers, Jan-2025 complaints (n=360)",
          {"size": 9, "italic": True, "color": GREEN_DK})],
    ], anchor=MSO_ANCHOR.MIDDLE, margins=0.18)
    # Right: the stat stack
    rx = Inches(8.2)
    rw = Inches(4.55)
    stats = [
        ("65%", "negative pharmacy sentiment", "vs 65% POSITIVE on retail price/promo — the trust gap is specific to pharmacy"),
        ("5", "complaint clusters", "chronic re-order (dominant) · delivery fails · accuracy · service/ownership · billing"),
        ("Defection", "to local pharmacies", "customers “willing to pay R460 more” for relationship-based care"),
    ]
    yy = top
    for big, mid, sub in stats:
        card = add_rect(slide, rx, yy, rw, Inches(1.42), fill=WHITE, line=LINE,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        add_rect(slide, rx, yy, Inches(0.1), Inches(1.42), fill=GREEN)
        add_text(slide, rx + Inches(0.28), yy + Inches(0.16), rw - Inches(0.5), Inches(0.5),
                 [[(big + "  ", {"size": 24, "bold": True, "color": GREEN}),
                   (mid, {"size": 13, "bold": True, "color": INK})]])
        add_text(slide, rx + Inches(0.28), yy + Inches(0.78), rw - Inches(0.5), Inches(0.55),
                 [[(sub, {"size": 9.5, "color": SLATE})]], line_spacing=1.0)
        yy += Inches(1.56)
    notes(slide, "Start with the problem, in the customers' own words. Social "
                 "listening on 1,270 coded comments shows pharmacy sentiment at 65% "
                 "negative — the mirror image of the retail brand, which is 65% "
                 "positive on price and promotions. So this is not a price or "
                 "availability problem; it is a trust and fulfilment problem. The "
                 "Patient Trust Journey shows how trust degrades: Calm -> Trigger -> "
                 "Acceleration -> Agitation -> Defection. The Jan-2025 complaints "
                 "(n=360) cluster into five themes, with chronic re-order dominant. "
                 "The endgame is defection: customers move scripts to local "
                 "pharmacies and pay more for relationship care. Note the honest "
                 "caveat: these are directional signals, not satisfaction rates.")
    return slide


def slide_approach(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Our approach", "How we built a defensible North Star",
        "A ground-truth research vault feeds a disciplined pipeline — research → "
        "best practice → principles → people → journey — so every design "
        "choice is traceable to evidence.",
        cite="CLAUDE.md (the vault schema) · meta/research-log.md")
    # The pipeline arrow row
    # Labels are kept short and rendered horizontally inside each chevron via a
    # dedicated full-width centred textbox (mirrors slides 3 & 13), so they never
    # stack one-letter-per-line or overflow the chevron.
    steps = [
        (["Research"], "8 domains,\n150+ sources"),
        (["Best", "practice"], "Global\nbenchmarks"),
        (["Strategic", "principles"], "S1–S5\nbusiness bets"),
        (["UX", "principles"], "U1–U7\nexperience rules"),
        (["Segments"], "Two-economy\ndual journey"),
        (["Personas"], "8 research-\nbacked people"),
        (["Journey"], "The North\nStar"),
    ]
    n = len(steps)
    top = Inches(2.85)
    aw = Inches(1.72)
    ah = Inches(1.5)
    overlap = Inches(0.18)
    total = aw * n - overlap * (n - 1)
    ax = (SW - total) / 2
    for i, (label_lines, sub) in enumerate(steps):
        x = ax + i * (aw - overlap)
        is_last = (i == n - 1)
        col = GREEN if is_last else (MINT if i % 2 == 0 else MINT_MD)
        chev = add_rect(slide, x, top, aw, ah, fill=col, shape=MSO_SHAPE.CHEVRON,
                        shadow=is_last)
        lab_color = WHITE if is_last else INK
        # Horizontal, centred label in its own textbox covering the chevron's
        # straight body (inset from the pointed ends so text stays contained).
        lab_x = x + Inches(0.26)
        lab_w = aw - Inches(0.52)
        add_text(slide, lab_x, top, lab_w, ah,
                 [[(ln, {"size": 10.5, "bold": True, "color": lab_color})]
                  for ln in label_lines],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=0.98)
        add_text(slide, x + Inches(0.1), top + ah + Inches(0.1), aw - Inches(0.2),
                 Inches(0.7),
                 [[(ln, {"size": 8.5, "color": SLATE})]
                  for ln in sub.split("\n")],
                 align=PP_ALIGN.CENTER, line_spacing=0.95)
    # The "trust the method" frame band
    by = Inches(5.35)
    band = add_rect(slide, Inches(0.6), by, Inches(12.13), Inches(1.35),
                    fill=CARD_BG, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(band, [
        [("The ground-truth vault method.  ", {"size": 13, "bold": True, "color": GREEN_DK}),
         ("Every claim traces to a cited source; synthesis lives in a structured "
          "wiki; design outputs are defensible, not opinion. ",
          {"size": 12, "color": INK})],
        [("This is the consulting discipline behind the North Star: "
          "we earn the outcome by trusting the method.",
          {"size": 12, "italic": True, "color": SLATE})],
    ], anchor=MSO_ANCHOR.MIDDLE, margins=0.2)
    notes(slide, "This is the 'trust the method' slide — the spine of a consulting "
                 "narrative. We did not start with a concept; we built a ground-truth "
                 "knowledge vault and ran a disciplined pipeline: research across 8 "
                 "domains, distil global best practice, derive strategic principles "
                 "(S1-S5) and UX principles (U1-U7), ground in the two-economy "
                 "segmentation and 8 personas, and only THEN design the journey. Every "
                 "downstream choice is traceable to evidence. Walk left to right — each "
                 "of the next sections is one chevron in this row.")
    return slide


def slide_evidence(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "The evidence base", "Scale and credibility behind the work",
        "~150+ cited sources across 8 domains and A–D reliability tiers, anchored "
        "by Dis-Chem's own first-party primary research — this is grounded, not "
        "borrowed opinion.",
        cite="sources/_sources-index.md · meta/research-log.md · meta/source-reliability-guide.md")
    top = Inches(2.4)
    # Left: four headline metric tiles (2x2)
    tiles = [
        ("8", "research domains", "regulatory, schemes, dispensing, digital, global, UX, consumer, Dis-Chem"),
        ("150+", "cited sources", "one note per source: citation, claims, verbatim excerpts"),
        ("A–D", "reliability tiers", "regulatory & clinical claims rest on A/B; signals labelled"),
        ("50+", "synthesis pages", "cross-linked wiki — claims link back to their evidence"),
    ]
    tw = Inches(3.55)
    th = Inches(1.85)
    gx = Inches(0.25)
    gy = Inches(0.22)
    lx = Inches(0.6)
    for idx, (big, mid, sub) in enumerate(tiles):
        r, c = divmod(idx, 2)
        x = lx + c * (tw + gx)
        y = top + r * (th + gy)
        add_rect(slide, x, y, tw, th, fill=CARD_BG, line=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        add_text(slide, x + Inches(0.25), y + Inches(0.18), tw - Inches(0.5), Inches(0.7),
                 [[(big, {"size": 34, "bold": True, "color": GREEN})]])
        add_text(slide, x + Inches(0.27), y + Inches(0.88), tw - Inches(0.5), Inches(0.3),
                 [[(mid, {"size": 13, "bold": True, "color": INK})]])
        add_text(slide, x + Inches(0.27), y + Inches(1.2), tw - Inches(0.5), Inches(0.6),
                 [[(sub, {"size": 9, "color": SLATE})]], line_spacing=0.98)
    # Right: first-party primary research panel
    rx = Inches(8.25)
    rw = Inches(4.48)
    add_rect(slide, rx, top, rw, Inches(3.92), fill=GREEN,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_text(slide, rx + Inches(0.3), top + Inches(0.22), rw - Inches(0.6), Inches(0.6),
             [[("PLUS DIS-CHEM'S OWN", {"size": 11, "bold": True, "color": LIME})],
              [("PRIMARY RESEARCH", {"size": 11, "bold": True, "color": LIME})]],
             line_spacing=1.0)
    primary = [
        ("Complaints analysis", "n=360 resolved (Jan 2025)"),
        ("Social listening", "1,270 coded comments (Dec 2025)"),
        ("Script-journey interviews", "n=10 IDIs (Oct 2025)"),
        ("Store-of-the-Future interviews", "n=12 (Jan 2026)"),
        ("Omnichannel tech roadmap", "24 capabilities (internal board)"),
    ]
    yy = top + Inches(1.0)
    for title_, meta in primary:
        add_text(slide, rx + Inches(0.3), yy, rw - Inches(0.6), Inches(0.6),
                 [[("▸ ", {"size": 12, "bold": True, "color": LIME}),
                   (title_, {"size": 12.5, "bold": True, "color": WHITE}),
                   ("   " + meta, {"size": 10, "color": MINT_MD})]],
                 line_spacing=1.0)
        yy += Inches(0.57)
    notes(slide, "Credibility. The North Star sits on roughly 150+ cited sources "
                 "spanning 8 research domains, each rated on an A-D reliability scale "
                 "so the weight of every claim is explicit — regulatory and clinical "
                 "claims rest on A/B sources, and consumer sentiment is labelled as a "
                 "signal. Crucially, this is not just desk research: it is anchored by "
                 "Dis-Chem's own first-party primary research — the complaints data, "
                 "social listening, two interview studies, and the internal "
                 "omnichannel tech roadmap. That mix of external rigour and internal "
                 "voice is what makes the North Star defensible to this room.")
    return slide


def slide_best_practice(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Best practice", "What the world's leading pharmacies prove — and what transfers to SA",
        "Borrow mechanisms, not funding architectures: Australia, India and China "
        "transfer to South Africa far better than the US, UK or EU models.",
        cite="wiki/global-benchmarks/patterns-to-borrow-and-cautions")
    top = Inches(2.4)
    # Left: patterns to borrow
    add_text(slide, Inches(0.6), top, Inches(7.0), Inches(0.32),
             [[("PATTERNS TO BORROW", {"size": 11.5, "bold": True, "color": GREEN})]])
    patterns = [
        ("QR e-script token", "scanned at any pharmacy; printed/USSD/WhatsApp fallback (AU)"),
        ("Nominated pharmacy / script list", "patient-consented, revocable (UK EPS, AU ASL)"),
        ("Auto-refill / med-sync", "the refill is automatic; cancelling is the exception"),
        ("Teleconsult → Rx → deliver", "one integrated loop, scoped to follow-up/chronic (IN, CN)"),
        ("Adherence packaging + app wallet", "dose-by-time packs; carer mode (PillPack, MedAdvisor)"),
        ("Price legible at point of commit", "SEP + capped fee already make this natural in SA"),
    ]
    yy = top + Inches(0.42)
    for name, sub in patterns:
        add_text(slide, Inches(0.6), yy, Inches(7.1), Inches(0.5),
                 [[("✓  ", {"size": 12, "bold": True, "color": GREEN}),
                   (name, {"size": 12, "bold": True, "color": INK}),
                   ("  —  " + sub, {"size": 10.5, "color": SLATE})]],
                 line_spacing=1.0)
        yy += Inches(0.62)
    # Right: transfer-strength table
    rx = Inches(8.05)
    rw = Inches(4.68)
    add_text(slide, rx, top, rw, Inches(0.32),
             [[("TRANSFER STRENGTH TO SA", {"size": 11.5, "bold": True, "color": GREEN})]])
    rows = [
        ("Australia", "Highest", GREEN),
        ("India", "High", GREEN),
        ("China", "High (loop + guardrails)", GREEN),
        ("UK", "Medium", RGBColor(0xC9, 0xA2, 0x2E)),
        ("EU / Germany", "Medium (as a warning)", RGBColor(0xC9, 0xA2, 0x2E)),
        ("USA", "Lowest (features only)", RGBColor(0xC9, 0x6A, 0x2E)),
    ]
    yy = top + Inches(0.42)
    for region, strength, col in rows:
        rcard = add_rect(slide, rx, yy, rw, Inches(0.52), fill=WHITE, line=LINE,
                         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(slide, rx, yy, Inches(0.1), Inches(0.52), fill=col)
        add_text(slide, rx + Inches(0.3), yy + Inches(0.07), Inches(2.0), Inches(0.4),
                 [[(region, {"size": 12, "bold": True, "color": INK})]],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, rx + Inches(2.2), yy + Inches(0.07), rw - Inches(2.4), Inches(0.4),
                 [[(strength, {"size": 11, "color": col})]], anchor=MSO_ANCHOR.MIDDLE)
        yy += Inches(0.58)
    add_text(slide, rx, yy + Inches(0.02), rw, Inches(0.5),
             [[("Reject: the US PBM layer, a single online-only rail, PIN/smart-card "
                "redemption, density-dependent speed promises.",
                {"size": 8.5, "italic": True, "color": SLATE})]], line_spacing=0.95)
    notes(slide, "We did not invent from scratch — we studied the world's leading "
                 "pharmacy models and asked what actually transfers to South Africa. "
                 "The headline: emerging-market and Australian models fit SA far "
                 "better than US/UK/EU, because they assume mixed funding, cash-paying "
                 "patients and dense store estates as fulfilment nodes — exactly SA's "
                 "reality. Borrow the mechanisms on the left (token e-scripts, "
                 "nominated pharmacy, auto-refill, teleconsult loops, adherence "
                 "packaging, price transparency); reject the funding architectures on "
                 "the right (especially the US PBM layer and Germany's single online-"
                 "only rail, which drove 99% of patients to the paper fallback).")
    return slide


def slide_strategic(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Strategic principles", "S1–S5 — the business bets",
        "Five strategic bets define what business this is: compete on service (not "
        "regulated price), own the chronic flywheel, and serve both economies on one "
        "spine.",
        cite="deliverables/_shared/pillars.md")
    bets = [
        ("S1", "Serve both economies on one spine",
         "Insured (~16%) and cash (~84%) are different journeys sharing one spine; the cash majority is never a degraded afterthought."),
        ("S2", "Pharmacy as the health front door",
         "Scripting is the entry to an ongoing care relationship, not a one-off transaction; the store + clinic estate is a strategic asset."),
        ("S3", "Compete on service, not price",
         "SEP + capped dispensing fee fix the price; the durable differentiator is the experience around the drug."),
        ("S4", "Chronic-medicine flywheel = retention engine",
         "Proactive, managed chronic repeats are the recurring relationship that compounds loyalty and adherence."),
        ("S5", "Pharmacist at top-of-licence; automate the toil",
         "Automation, central-fill and telepharmacy absorb the mechanics so pharmacists do the clinical, trust-building work."),
    ]
    _pillar_grid(slide, bets, top=Inches(2.45), cols=3, accent=GREEN)
    notes(slide, "The strategic layer answers: why build this, and what business is "
                 "it? Five bets. S1: design for the two South Africas at once. S2: "
                 "treat the pharmacy as a health front door, not a counter. S3 is the "
                 "keystone — because SEP and the capped dispensing fee fix the medicine "
                 "price, you literally cannot compete on price; the only durable "
                 "differentiator is the experience. S4: the chronic flywheel is the "
                 "retention engine the courier pharmacies currently win. S5: automate "
                 "the toil so pharmacists work at the top of their licence. Every step "
                 "of the journey is defensible from at least one of these.")
    return slide


def slide_ux(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "UX principles", "U1–U7 — the experience constitution",
        "Seven UX rules govern how the experience must behave at every touchpoint — "
        "WhatsApp-first, never leave them guessing, money legible, graceful failure, "
        "dignity, fewer steps, adherence by design.",
        cite="deliverables/_shared/pillars.md")
    rules = [
        ("U1", "Channel ladder: WhatsApp-first & data-frugal",
         "WhatsApp/USSD/SMS first, app as enhancement; low-data, async, resumable."),
        ("U2", "Never leave them guessing",
         "Honest real-time status at every step; the silent wait is the primary pain."),
        ("U3", "Make money legible, early, always",
         "Price + cover + co-pay shown before the till; cash/EFT/wallet; no cost shock."),
        ("U4", "Design for graceful failure",
         "Plain-language recovery, never a dead end, always a route to a pharmacist."),
        ("U5", "Dignity, privacy & language by default",
         "Stigma-aware, discreet, multilingual; caregivers/proxies are first-class users."),
        ("U6", "Reduce the journey, don't just digitise it",
         "Eliminate steps and re-keying; one profile across channels; measure friction removed."),
        ("U7", "Adherence-as-design",
         "Design the schedule, label and refill loop as a behaviour-change surface."),
    ]
    _pillar_grid(slide, rules, top=Inches(2.4), cols=4, accent=GREEN_DK,
                 card_h=Inches(2.02), small=True)
    notes(slide, "If the strategic bets are the 'what business', the UX principles are "
                 "the 'how it must behave' — the experience constitution, checkable at "
                 "every touchpoint. U1 meets people on WhatsApp, the channel they "
                 "already have. U2 — never leave them guessing — directly answers the "
                 "silent backstage wait that drives the trust journey. U3 makes money "
                 "legible before commitment. U4 designs for graceful failure with a "
                 "human fallback. U5 protects dignity and treats caregivers as first-"
                 "class. U6 removes steps rather than prettifying them. U7 designs for "
                 "adherence, not just dispensing. A good design step satisfies at least "
                 "one S and one U.")
    return slide


def _pillar_grid(slide, items, top, cols, accent, card_h=None, small=False):
    n = len(items)
    rows = (n + cols - 1) // cols
    gx = Inches(0.22)
    gy = Inches(0.22)
    lx = Inches(0.6)
    avail_w = SW - lx * 2
    cw = (avail_w - gx * (cols - 1)) / cols
    if card_h is None:
        card_h = Inches(1.72)
    for idx, (pid, name, desc) in enumerate(items):
        r, c = divmod(idx, cols)
        x = lx + c * (cw + gx)
        y = top + r * (card_h + gy)
        card = add_rect(slide, x, y, cw, card_h, fill=CARD_BG, line=LINE,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        # id badge
        badge = add_rect(slide, x + Inches(0.18), y + Inches(0.18), Inches(0.62),
                         Inches(0.42), fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(badge, [[(pid, {"size": 14, "bold": True, "color": WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
        nsize = 11.5 if not small else 11
        dsize = 9.5 if not small else 8.7
        add_text(slide, x + Inches(0.95), y + Inches(0.16), cw - Inches(1.1), Inches(0.7),
                 [[(name, {"size": nsize, "bold": True, "color": INK})]],
                 line_spacing=0.95)
        add_text(slide, x + Inches(0.2), y + Inches(0.78 if not small else 0.74),
                 cw - Inches(0.4), card_h - Inches(0.9),
                 [[(desc, {"size": dsize, "color": SLATE})]], line_spacing=1.0)


def slide_segments(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Segments — the two-economy reality", "One spine, two first-class funded paths",
        "Funding — not age, device or geography — is the single most load-bearing "
        "design variable; the cash-paying majority must never feel like a degraded "
        "afterthought.",
        cite="wiki/concepts/dual-journey · wiki/consumer-context/two-economy-affordability")
    top = Inches(2.45)
    # split bar: 16% insured / 84% cash
    barx = Inches(0.6)
    barw = Inches(12.13)
    barh = Inches(0.82)
    insured_w = Emu(int(barw * 0.16))
    cash_w = Emu(int(barw) - int(insured_w))
    ins = add_rect(slide, barx, top, insured_w, barh, fill=GREEN,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(ins, [[("16%", {"size": 15, "bold": True, "color": WHITE})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
    cash = add_rect(slide, Emu(int(barx) + int(insured_w)), top, cash_w, barh,
                    fill=MINT_MD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(cash, [[("84%", {"size": 15, "bold": True, "color": INK})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
    add_text(slide, barx, top + barh + Inches(0.04), Inches(3.0), Inches(0.3),
             [[("Insured (medical scheme)", {"size": 10, "bold": True, "color": GREEN_DK})]])
    add_text(slide, Emu(int(barx) + int(insured_w)) + Inches(0.1),
             top + barh + Inches(0.04), Inches(6.0), Inches(0.3),
             [[("Cash / uninsured majority", {"size": 10, "bold": True, "color": SLATE})]])
    # two columns comparing the journeys
    cy = top + Inches(1.55)
    colw = Inches(5.95)
    lx = Inches(0.6)
    rx = lx + colw + Inches(0.25)
    # insured column
    add_rect(slide, lx, cy, colw, Inches(2.7), fill=MINT, line=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_text(slide, lx + Inches(0.3), cy + Inches(0.2), colw - Inches(0.6), Inches(0.4),
             [[("INSURED JOURNEY", {"size": 12, "bold": True, "color": GREEN_DK}),
               ("   — Nomvula", {"size": 11, "italic": True, "color": SLATE})]])
    ins_rows = [
        "Fund: real-time claim pre-adjudication; co-pay shown before commit",
        "Fulfilment: home delivery or collection by default",
        "Pay: scheme + saved method",
        "Steering: DSP / formulary / co-pay logic",
    ]
    _bullet_block(slide, lx + Inches(0.3), cy + Inches(0.66), colw - Inches(0.6),
                  ins_rows, size=11, gap=0.5)
    # cash column
    add_rect(slide, rx, cy, colw, Inches(2.7), fill=CARD_BG, line=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_text(slide, rx + Inches(0.3), cy + Inches(0.2), colw - Inches(0.6), Inches(0.4),
             [[("CASH / UNINSURED JOURNEY", {"size": 12, "bold": True, "color": INK}),
               ("   — Sipho", {"size": 11, "italic": True, "color": SLATE})]])
    cash_rows = [
        "Fund: transparent cash price (SEP + dispensing fee) before commit",
        "Fulfilment: nearest collection / locker (delivery = premium)",
        "Pay: cash-on-collection + instant EFT / wallet, no saved card",
        "Steering: generic-first for affordability; CCMDD precedent",
    ]
    _bullet_block(slide, rx + Inches(0.3), cy + Inches(0.66), colw - Inches(0.6),
                  cash_rows, size=11, gap=0.5, marker_color=SLATE)
    notes(slide, "The most load-bearing design variable is not age or device — it is "
                 "funding. South Africa is a two-economy market: roughly 16% on a "
                 "medical scheme, 84% paying cash or relying on the public sector. The "
                 "design answer is one service spine with two first-class funded paths, "
                 "diverging only where money works differently — pre-adjudicated co-pay "
                 "on the insured side, transparent cash price on the cash side. The "
                 "non-negotiable: the cash majority is never a degraded afterthought. "
                 "We anchor the two paths on real personas — Nomvula (insured chronic) "
                 "and Sipho (cash, low-data). Removing a trip or a surprise charge is a "
                 "real financial benefit to the cash customer, not just convenience.")
    return slide


def slide_personas(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Personas", "Eight research-backed people — three become the protagonists",
        "The journey is designed for real, evidenced people spanning the two-economy "
        "split — spotlighting an insured chronic patient, a cash WhatsApp-first "
        "customer, and a caregiver.",
        cite="design-kit/personas/_personas-index")
    top = Inches(2.4)
    # spotlight three (top row, bigger cards), then a strip of the other five
    spot = [
        ("Nomvula Khumalo", "Insured chronic, multimorbid",
         "54, Discovery, HTN + type-2 diabetes on the CDL. Wants several repeats to flow automatically, fully funded, no co-pay surprise."),
        ("Sipho Ndlovu", "Cash, low-data, WhatsApp-first",
         "45, peri-urban, cheap Android, no street address. WhatsApp + SMS PIN + a collection point near home are the only things that reliably work."),
        ("Lerato Dlamini", "Caregiver / proxy collector",
         "38, sandwich-generation. Runs medicine admin for her diabetic mother, two children and herself — multiple profiles, cycles and wallets."),
    ]
    cw = Inches(3.95)
    gx = Inches(0.14)
    lx = Inches(0.6)
    for i, (name, role, desc) in enumerate(spot):
        x = lx + i * (cw + gx)
        add_rect(slide, x, top, cw, Inches(2.55), fill=CARD_BG, line=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        # avatar circle
        av = add_rect(slide, x + Inches(0.28), top + Inches(0.26), Inches(0.7),
                      Inches(0.7), fill=GREEN, shape=MSO_SHAPE.OVAL)
        initials = "".join(w[0] for w in name.split()[:2])
        shape_text(av, [[(initials, {"size": 16, "bold": True, "color": WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
        add_text(slide, x + Inches(1.1), top + Inches(0.28), cw - Inches(1.3), Inches(0.4),
                 [[(name, {"size": 13.5, "bold": True, "color": INK})]])
        add_text(slide, x + Inches(1.1), top + Inches(0.66), cw - Inches(1.3), Inches(0.3),
                 [[(role, {"size": 9.5, "bold": True, "color": GREEN_DK})]])
        add_text(slide, x + Inches(0.28), top + Inches(1.18), cw - Inches(0.56), Inches(1.25),
                 [[(desc, {"size": 10, "color": SLATE})]], line_spacing=1.02)
        # spotlight tab
        tab = add_rect(slide, x + cw - Inches(1.35), top + Inches(0.1), Inches(1.2),
                       Inches(0.3), fill=LIME, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(tab, [[("★ PROTAGONIST", {"size": 7.5, "bold": True, "color": INK})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
    # the other five as a slim strip
    others = [
        "Thabo — cash-pay acute",
        "Trudie — elderly poly-pharmacy",
        "Aisha — time-poor professional",
        "Mandla — public→retail bridge",
        "Naledi — frontline pharmacist (staff)",
    ]
    sy = top + Inches(2.78)
    add_text(slide, lx, sy, Inches(6.0), Inches(0.3),
             [[("THE FULL SET ALSO SPANS:", {"size": 10, "bold": True, "color": GREEN})]])
    ow = Inches(2.34)
    ogx = Inches(0.12)
    for i, o in enumerate(others):
        x = lx + i * (ow + ogx)
        chip = add_rect(slide, x, sy + Inches(0.34), ow, Inches(0.62), fill=MINT,
                        line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(chip, [[(o, {"size": 9.5, "bold": True, "color": INK})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.06)
    notes(slide, "The journey is grounded in eight research-backed personas — not "
                 "invented archetypes; every SA-context claim cites a wiki page. They "
                 "deliberately span the two-economy split, the digital-literacy and "
                 "connectivity spread, and the customer-vs-staff tension. Three become "
                 "the protagonists of the North Star: Nomvula, the insured chronic "
                 "multimorbid patient; Sipho, the cash, low-data, WhatsApp-first "
                 "customer who has no street address; and Lerato, the caregiver "
                 "managing meds for three generations. The remaining five — including "
                 "Naledi, the frontline pharmacist whose oversight we must protect — "
                 "keep the design honest across the rest of the market.")
    return slide


def slide_pains_mapped(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "The pains, mapped", "Today's complaint clusters map cleanly onto the journey spine",
        "Most root causes sit in the digital and fulfilment seams — yet surface as "
        "'pharmacy' complaints; fixing the counter alone won't fix them. This sets up "
        "the journey as the answer.",
        cite="wiki/dischem/pharmacy-complaints-pain-points · src-dc-complaints-jan2025")
    top = Inches(2.4)
    # mini spine
    spine = ["SUBMIT", "VALIDATE", "FUND / PAY", "DISPENSE", "DELIVER", "REFILL"]
    sx = Inches(0.6)
    sw_ = Inches(2.0)
    for i, st in enumerate(spine):
        x = sx + i * (sw_ - Inches(0.04))
        ch = add_rect(slide, x, top, sw_, Inches(0.5), fill=GREEN if i % 2 == 0 else GREEN_DK,
                      shape=MSO_SHAPE.CHEVRON)
        shape_text(ch, [[(st, {"size": 10, "bold": True, "color": WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
    # five clusters as cards mapped to stages, each with a quote
    clusters = [
        ("A · Chronic re-order", "REFILL / SUBMIT",
         "“Dis-Chem App is dysfunctional.”", "Dominant cluster"),
        ("B · Delivery fails", "DELIVER",
         "“Out for delivery. I did not receive it.”", "Priority 2"),
        ("C · Dispensing accuracy", "DISPENSE",
         "“My blood-pressure tablets were missing.”", "Higher-risk"),
        ("D · Service & ownership", "Cross-cutting",
         "“The nurse doesn't know what to do.”", "Priority 4"),
        ("E · Billing transparency", "FUND / PAY",
         "“Pricing differs dramatically between stores.”", "Priority 5"),
    ]
    cy = top + Inches(0.95)
    cw = Inches(2.37)
    gx = Inches(0.1)
    for i, (name, stage, quote, tag) in enumerate(clusters):
        x = Inches(0.6) + i * (cw + gx)
        add_rect(slide, x, cy, cw, Inches(3.1), fill=CARD_BG, line=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        add_rect(slide, x, cy, cw, Inches(0.06), fill=GREEN)
        add_text(slide, x + Inches(0.16), cy + Inches(0.16), cw - Inches(0.32), Inches(0.55),
                 [[(name, {"size": 11, "bold": True, "color": INK})]], line_spacing=0.95)
        tagb = add_rect(slide, x + Inches(0.16), cy + Inches(0.72), Inches(1.7), Inches(0.28),
                        fill=MINT_MD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(tagb, [[(tag, {"size": 8, "bold": True, "color": GREEN_DK})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
        add_text(slide, x + Inches(0.16), cy + Inches(1.15), cw - Inches(0.32), Inches(0.3),
                 [[("→ " + stage, {"size": 8.5, "bold": True, "color": GREEN})]])
        add_text(slide, x + Inches(0.16), cy + Inches(1.55), cw - Inches(0.32), Inches(1.4),
                 [[(quote, {"size": 9.5, "italic": True, "color": SLATE})]],
                 line_spacing=1.05)
    notes(slide, "Here is the bridge from problem to solution. The five complaint "
                 "clusters map cleanly onto the journey spine: chronic re-order "
                 "(dominant) at REFILL/SUBMIT; delivery failure at DELIVER; dispensing "
                 "accuracy at DISPENSE; service and ownership cross-cutting; and "
                 "billing transparency at FUND/PAY. The load-bearing insight: most root "
                 "causes actually live in the digital and fulfilment seams, but they "
                 "surface as pharmacy complaints because the pharmacy is the visible "
                 "face. So designing only the counter won't fix them — we have to fix "
                 "the seams. That is exactly what the North Star journey does, which is "
                 "where we go next.")
    return slide


def slide_northstar_intro(prs):
    """Reveal slide 1: the concept + board screenshot."""
    slide = new_slide(prs)
    content_chrome(
        slide, "The North Star — reveal", "The household medicine manager — a self-running repeat",
        "One person manages her own, her mother's and her child's medication in one "
        "place; the system predicts, dispenses and delivers — she barely has to think "
        "about it.",
        cite="deliverables/05-tobe-future-state-journey.html · src-vn-cape-point-11")
    # left text, right board image
    lx = Inches(0.6)
    add_text(slide, lx, Inches(2.5), Inches(4.3), Inches(0.4),
             [[("THE MOMENT ARC", {"size": 11, "bold": True, "color": GREEN})]])
    moments = [
        "“My repeat is due — and I barely have to think about it”",
        "“I handle Mum's and my little one's meds in one place”",
        "“I know the cost and cover before I commit”",
        "“A named pharmacist has my back”",
        "“It's picked, checked and packed — and I watch it happen”",
        "“It arrives when they said — never left guessing”",
        "“When life throws a curveball, the system bends”",
    ]
    yy = Inches(2.92)
    for m in moments:
        add_text(slide, lx, yy, Inches(4.35), Inches(0.5),
                 [[("→ ", {"size": 10, "bold": True, "color": GREEN}),
                   (m, {"size": 10.5, "italic": True, "color": INK})]],
                 line_spacing=0.98)
        yy += Inches(0.5)
    # board image, framed, right
    iw = Inches(7.4)
    ih = Emu(int(iw / 1.524))
    ix = Inches(5.25)
    iy = Inches(2.55)
    add_rect(slide, ix - Inches(0.06), iy - Inches(0.06), iw + Inches(0.12),
             ih + Inches(0.12), fill=WHITE, line=LINE, line_w=Pt(1), shadow=True)
    slide.shapes.add_picture(IMG_BOARD, ix, iy, width=iw, height=ih)
    add_text(slide, ix, iy + ih + Inches(0.06), iw, Inches(0.3),
             [[("The built artefact — interactive future-state journey board "
                "(insured view)", {"size": 8.5, "italic": True, "color": SLATE})]],
             align=PP_ALIGN.CENTER)
    notes(slide, "This is the climax — the North Star itself, and the actual built "
                 "artefact, not a sketch. The complex use case is the household "
                 "medicine manager: one person, simultaneously patient and caregiver, "
                 "managing her own multimorbid chronic meds, her elderly mother's "
                 "chronic meds, and her child's occasional acute needs — across one "
                 "account. It exercises the whole spine. The moment arc on the left is "
                 "written in the customer's first-person voice; each moment answers a "
                 "real, quoted pain. This is an aspirational ~5-year, tech-led vision, "
                 "deliberately unconstrained by today's technology landscape.")
    return slide


def slide_northstar_moment(prs):
    """Reveal slide 2: expanded moment + cash view side by side."""
    slide = new_slide(prs)
    content_chrome(
        slide, "The North Star — reveal", "Calm on the surface, depth on demand — and equal on both paths",
        "Each moment opens to reveal the experience behind the line of visibility; "
        "the dual-economy toggle proves the cash path is first-class, not a "
        "downgrade.",
        cite="deliverables/05-tobe-future-state-journey.html · wiki/concepts/dual-journey")
    top = Inches(2.5)
    # two images side by side (both landscape)
    iw = Inches(6.0)
    ih = Emu(int(iw / 1.524))
    gap = Inches(0.3)
    lx = Inches(0.6)
    rx = lx + iw + gap
    for img, x, cap in [
        (IMG_MOMENT, lx, "A moment expanded — story + behind-the-line (insured)"),
        (IMG_CASH, rx, "The same journey, cash / uninsured path (Sipho)"),
    ]:
        add_rect(slide, x - Inches(0.05), top - Inches(0.05), iw + Inches(0.1),
                 ih + Inches(0.1), fill=WHITE, line=LINE, line_w=Pt(1), shadow=True)
        slide.shapes.add_picture(img, x, top, width=iw, height=ih)
        add_text(slide, x, top + ih + Inches(0.06), iw, Inches(0.3),
                 [[(cap, {"size": 9, "italic": True, "color": SLATE})]],
                 align=PP_ALIGN.CENTER)
    # bottom takeaways
    by = top + ih + Inches(0.5)
    band = add_rect(slide, Inches(0.6), by, Inches(12.13), Inches(0.72), fill=MINT,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(band, [[
        ("Low-friction & human:  ", {"size": 11.5, "bold": True, "color": GREEN_DK}),
        ("predictability beats speed (SOTF), the customer surface stays calm, and the "
         "ops/tech depth lives behind a line of visibility the reader can open. The "
         "cash majority gets the same dignity, transparency and fulfilment quality.",
         {"size": 11, "color": INK})]],
        anchor=MSO_ANCHOR.MIDDLE, margins=0.2)
    notes(slide, "Two things to land here. First, the experience is calm on the "
                 "surface but deep on demand: a customer sees a confident, human, "
                 "first-person story; a reviewer can open each moment to see the "
                 "frontstage, line of visibility, backstage, tech callouts, pains and "
                 "opportunities behind it. Predictability beats speed — that is the "
                 "Store-of-the-Future finding. Second, the dual-economy toggle: the "
                 "left is the insured path, the right is the same journey for the cash, "
                 "low-data customer. The cash path is genuinely first-class — cash "
                 "price legible up front, locker/collection options, no scheme to lean "
                 "on — not a degraded afterthought.")
    return slide


def slide_northstar_cabinet(prs):
    """Reveal slide 3: the medicine cabinet hero moment (portrait image)."""
    slide = new_slide(prs)
    content_chrome(
        slide, "The North Star — reveal", "“My medicine cabinet” — the hero moment",
        "Every household member's medication in one beautiful view — what to take, "
        "when, what's left and what's coming — turns admin into confidence.",
        cite="deliverables/05-tobe-future-state-journey.html · src-vn-cape-point-11")
    # portrait image centered-right
    ih = Inches(4.35)
    iw = Emu(int(ih * 0.647))
    ix = Inches(8.55)
    iy = Inches(2.45)
    add_rect(slide, ix - Inches(0.05), iy - Inches(0.05), iw + Inches(0.1),
             ih + Inches(0.1), fill=WHITE, line=LINE, line_w=Pt(1), shadow=True)
    slide.shapes.add_picture(IMG_CABINET, ix, iy, height=ih)
    add_text(slide, ix - Inches(0.5), iy + ih + Inches(0.08), iw + Inches(1.0), Inches(0.3),
             [[("★ The “manage my medications” showpiece moment",
                {"size": 8.5, "italic": True, "color": SLATE})]],
             align=PP_ALIGN.CENTER)
    # left: what it does
    lx = Inches(0.6)
    add_text(slide, lx, Inches(2.5), Inches(7.4), Inches(0.4),
             [[("WHY IT MATTERS", {"size": 11, "bold": True, "color": GREEN})]])
    points = [
        ("Every member, one view", "her own multimorbid regimen, her mother's chronic meds, her child's acute scripts — with a clear pill visual."),
        ("Predict, don't prompt", "the cabinet shows what's running low and pre-stages the repeat — the best refill is the one she didn't have to think about (U7, S4)."),
        ("Adherence by design", "what to take and when, with hybrid digital + printed instructions and consented, opt-out-able reminders tied to routine."),
        ("Dignity & proxy by default", "caregivers are first-class users; POPIA-grade consent for managing a dependant's medication (U5)."),
    ]
    yy = Inches(2.95)
    for lead, rest in points:
        add_text(slide, lx, yy, Inches(7.55), Inches(0.85),
                 [[("▸ ", {"size": 12, "bold": True, "color": GREEN}),
                   (lead, {"size": 12.5, "bold": True, "color": INK}),
                   ("  —  " + rest, {"size": 11, "color": SLATE})]],
                 line_spacing=1.05)
        yy += Inches(0.98)
    notes(slide, "The hero moment Jian and Tamsin called out specifically: a "
                 "'really beautiful end-to-end' view of every household member's "
                 "medication, with a pill visual, tracking, and 'using my medication'. "
                 "This is where the experience stops being a transaction and becomes a "
                 "relationship. It directly delivers the chronic flywheel (S4) and "
                 "adherence-as-design (U7): the cabinet predicts and pre-stages the "
                 "repeat rather than waiting for the customer to remember. And it treats "
                 "the caregiver as a first-class user with proper proxy consent — which "
                 "is the lived reality for millions of South African households.")
    return slide


def slide_makes_it_real(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "What makes it real", "The capability roadmap already exists — and we flag what needs reform",
        "ROWA, OCR, AI capture, BRIX, the E2E Dashboard and PICCUP are Dis-Chem's own "
        "roadmapped capabilities; the vision is buildable, with reform-dependencies "
        "named honestly.",
        cite="wiki/digital-transformation/dischem-scripting-tech-roadmap · src-dc-omnichannel-scripting-roadmap")
    top = Inches(2.4)
    # left: capability chips
    add_text(slide, Inches(0.6), top, Inches(7.0), Inches(0.32),
             [[("ROADMAPPED CAPABILITIES (planned, not yet live)",
                {"size": 11, "bold": True, "color": GREEN})]])
    caps = [
        ("ROWA", "automated picker / dispensing robot — the proof-point for 'automate the toil' (S5)"),
        ("OCR + AI capture", "reads, standardises & routes the script; human fallback (PBQ)"),
        ("BRIX", "auto-dispense of eligible repeats + point-of-sale"),
        ("E2E Dashboard", "one workflow manager across dispensary, store, collection & delivery"),
        ("PICCUP + VC", "consolidation → driver → last-mile with proof-of-delivery"),
        ("Virtual Manager", "human SLA oversight of the virtual order queue"),
    ]
    yy = top + Inches(0.45)
    for name, desc in caps:
        chip = add_rect(slide, Inches(0.6), yy, Inches(1.85), Inches(0.5), fill=GREEN,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(chip, [[(name, {"size": 11.5, "bold": True, "color": WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.04)
        add_text(slide, Inches(2.6), yy + Inches(0.05), Inches(5.05), Inches(0.5),
                 [[(desc, {"size": 10, "color": SLATE})]], line_spacing=0.98,
                 )
        yy += Inches(0.62)
    # right: reform register (amber)
    rx = Inches(8.0)
    rw = Inches(4.73)
    add_rect(slide, rx, top, rw, Inches(3.95), fill=RGBColor(0xFD, 0xF6, 0xE8),
             line=RGBColor(0xE6, 0xCB, 0x8A), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, rx + Inches(0.28), top + Inches(0.2), rw - Inches(0.56), Inches(0.4),
             [[("\U0001f52e  REQUIRES REFORM — flagged honestly",
                {"size": 11.5, "bold": True, "color": AMBER})]])
    reforms = [
        "Auto-dispense without human intervention vs pharmacist-oversight & valid-prescription rules",
        "Fully-paperless AES e-script acceptance (ECTA / Reg 33)",
        "Locker / unattended dispensing of scheduled medicines",
        "After-hours live scheme authorisation & 24/7 pharmacist oversight",
        "Remote pharmacist supervision / central-fill",
    ]
    yy = top + Inches(0.72)
    for r in reforms:
        add_text(slide, rx + Inches(0.28), yy, rw - Inches(0.56), Inches(0.65),
                 [[("▸ ", {"size": 11, "bold": True, "color": AMBER}),
                   (r, {"size": 10.5, "color": INK})]], line_spacing=1.0)
        yy += Inches(0.62)
    notes(slide, "Now we ground the vision so it does not read as science fiction. "
                 "Every capability on the left is Dis-Chem's OWN roadmapped technology, "
                 "captured from the internal omnichannel board — ROWA the picker robot, "
                 "OCR and AI capture, BRIX auto-dispense, the end-to-end Dashboard, and "
                 "PICCUP delivery. ROWA and OCR are already signed off. Important "
                 "honesty: these are planned, not confirmed live. And on the right we "
                 "name the reform-dependencies up front — fully automated dispensing, "
                 "paperless AES e-scripts, locker dispensing of scheduled meds, after-"
                 "hours authorisation, and remote supervision. Naming them protects "
                 "credibility and sets the regulatory agenda.")
    return slide


def slide_before_after(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "From pain to North Star", "Every future moment answers a real, quoted pain of today",
        "The gap is undeniable when you put today's customer words next to the future "
        "experience — the North Star is a direct response to the evidence.",
        cite="wiki/dischem/pharmacy-complaints-pain-points · dischem-social-sentiment-trust-journey")
    top = Inches(2.4)
    # header row
    hx1 = Inches(0.6)
    wcol = Inches(6.0)
    hx2 = hx1 + wcol + Inches(0.13)
    th = add_rect(slide, hx1, top, wcol, Inches(0.46),
                  fill=RGBColor(0xF3, 0xE4, 0xE4), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(th, [[("TODAY  (the quoted pain)", {"size": 12, "bold": True,
                                                   "color": RGBColor(0x9B, 0x37, 0x37)})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    fh = add_rect(slide, hx2, top, wcol, Inches(0.46), fill=GREEN,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(fh, [[("NORTH STAR  (the moment that answers it)",
                      {"size": 12, "bold": True, "color": WHITE})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    pairs = [
        ("“Pulling my hair trying to order via the App or WhatsApp.”",
         "My repeat is predicted and one-tap — I barely have to think about it."),
        ("“Out for delivery. I did not receive it.”",
         "Realistic windows, closed-loop tracking & proof-of-delivery — I'm never left guessing."),
        ("“My blood-pressure tablets were missing.”",
         "ROWA picks, the pharmacist signs off, an itemised handover I can check."),
        ("“The nurse doesn't know what to do.”",
         "A named pharmacist has my back; renewing is a tap, not a trek."),
        ("“Pricing differs dramatically between stores.”",
         "Cost and cover are legible before I commit — no surprise at the till."),
    ]
    ry = top + Inches(0.58)
    rh = Inches(0.74)
    for i, (pain, fut) in enumerate(pairs):
        bg = WHITE if i % 2 == 0 else CARD_BG
        add_rect(slide, hx1, ry, wcol, rh, fill=RGBColor(0xFB, 0xF3, 0xF3), line=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, hx1 + Inches(0.2), ry + Inches(0.08), wcol - Inches(0.4),
                 rh - Inches(0.16),
                 [[(pain, {"size": 11, "italic": True, "color": RGBColor(0x8A, 0x33, 0x33)})]],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.98)
        add_rect(slide, hx2, ry, wcol, rh, fill=MINT, line=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, hx2 + Inches(0.2), ry + Inches(0.08), wcol - Inches(0.4),
                 rh - Inches(0.16),
                 [[(fut, {"size": 11, "color": INK})]],
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.98)
        # arrow between
        ar = add_rect(slide, hx2 - Inches(0.16), ry + rh/2 - Inches(0.11),
                      Inches(0.22), Inches(0.22), fill=GREEN, shape=MSO_SHAPE.CHEVRON)
        ry += rh + Inches(0.1)
    notes(slide, "This is the slide that closes the loop on the whole narrative. On "
                 "the left, in red, are real customer words from the complaints and "
                 "social-listening data. On the right, in green, is the specific future "
                 "moment that answers each one. Pulling-my-hair re-orders become a "
                 "predicted one-tap repeat. 'I did not receive it' becomes closed-loop "
                 "tracking with proof-of-delivery. Missing tablets become ROWA-picked, "
                 "pharmacist-checked, itemised handover. The gap is undeniable — and it "
                 "shows the North Star is evidence-driven, not aspirational decoration.")
    return slide


def slide_whats_next(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "What's next", "This deck is the groundwork — here is the road to build",
        "We've defined the experience and earned buy-in; the next steps turn it into "
        "a service blueprint, per-channel adaptations, and a build.",
        cite="deliverables/05-tobe-future-state-journey.html (experience definition)")
    top = Inches(2.7)
    steps = [
        ("1", "Experience defined", "THIS DECK", "Agree the North Star — the experience we want to create.", True),
        ("2", "Service blueprint", "NEXT", "Frontstage/backstage lanes, ownership, SLAs, exception handling."),
        ("3", "Per-channel adaptation", "THEN", "Tune the base journey for WhatsApp, app, in-store and USSD."),
        ("4", "Build & pilot", "BUILD", "Sequence against the roadmap; pilot, measure, scale."),
    ]
    n = len(steps)
    cw = Inches(2.95)
    gx = Inches(0.13)
    lx = (SW - (cw * n + gx * (n - 1))) / 2
    for i, step in enumerate(steps):
        num, title_, tag, desc = step[0], step[1], step[2], step[3]
        here = len(step) > 4 and step[4]
        x = lx + i * (cw + gx)
        col = GREEN if here else CARD_BG
        card = add_rect(slide, x, top, cw, Inches(3.0), fill=col, line=LINE,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=here)
        circ = add_rect(slide, x + Inches(0.3), top + Inches(0.28), Inches(0.7),
                        Inches(0.7), fill=LIME if here else GREEN, shape=MSO_SHAPE.OVAL)
        shape_text(circ, [[(num, {"size": 20, "bold": True,
                                  "color": INK if here else WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
        tagb = add_rect(slide, x + cw - Inches(1.35), top + Inches(0.36), Inches(1.05),
                        Inches(0.34), fill=LIME if here else MINT_MD,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(tagb, [[(tag, {"size": 8.5, "bold": True, "color": INK})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
        add_text(slide, x + Inches(0.3), top + Inches(1.15), cw - Inches(0.6), Inches(0.6),
                 [[(title_, {"size": 15, "bold": True,
                             "color": WHITE if here else INK})]], line_spacing=0.95)
        add_text(slide, x + Inches(0.3), top + Inches(1.78), cw - Inches(0.6), Inches(1.1),
                 [[(desc, {"size": 10.5, "color": MINT_MD if here else SLATE})]],
                 line_spacing=1.05)
        # connector arrow
        if i < n - 1:
            add_rect(slide, x + cw - Inches(0.02), top + Inches(1.4), Inches(0.18),
                     Inches(0.22), fill=GREEN, shape=MSO_SHAPE.CHEVRON)
    notes(slide, "To close: be explicit about what this deck is and is not. This is "
                 "the experience definition — the groundwork that earns buy-in. It is "
                 "deliberately NOT yet a service blueprint. The road from here: (1) "
                 "agree the North Star today; (2) develop the full service blueprint "
                 "with frontstage/backstage lanes, ownership and SLAs; (3) adapt the "
                 "base journey per channel — WhatsApp, app, in-store, USSD; (4) sequence "
                 "the build against the roadmap, pilot, measure and scale. The ask of "
                 "this room is step one: agree the experience, so we can blueprint and "
                 "build it.")
    return slide


def slide_appendix_method(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Appendix", "Methodology & source-reliability tiers",
        "How the vault works and how we weight evidence — the discipline behind every "
        "cited claim in this deck.",
        cite="CLAUDE.md · meta/source-reliability-guide.md")
    top = Inches(2.4)
    # left: method
    lx = Inches(0.6)
    add_rect(slide, lx, top, Inches(6.0), Inches(3.9), fill=CARD_BG, line=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, lx + Inches(0.3), top + Inches(0.2), Inches(5.4), Inches(0.4),
             [[("THE GROUND-TRUTH VAULT (3 layers)", {"size": 11.5, "bold": True, "color": GREEN})]])
    method = [
        ("Raw sources", "one note per source: citation, key claims, verbatim excerpts — append-only evidence."),
        ("The wiki", "synthesised, cross-linked entity/concept pages — claims link back to their sources."),
        ("The schema", "conventions & workflows (ingest, query, lint) that keep it consistent and current."),
    ]
    yy = top + Inches(0.7)
    for name, desc in method:
        add_text(slide, lx + Inches(0.3), yy, Inches(5.4), Inches(0.85),
                 [[(name + ".  ", {"size": 11.5, "bold": True, "color": INK}),
                   (desc, {"size": 10.5, "color": SLATE})]], line_spacing=1.02)
        yy += Inches(0.92)
    add_text(slide, lx + Inches(0.3), yy + Inches(0.05), Inches(5.4), Inches(0.5),
             [[("Golden rule: ", {"size": 10.5, "bold": True, "color": GREEN_DK}),
               ("a wiki claim links to the source that backs it — synthesis, not narrative, in evidence.",
                {"size": 10.5, "italic": True, "color": SLATE})]], line_spacing=1.0)
    # right: tiers
    rx = Inches(6.95)
    rw = Inches(5.78)
    add_text(slide, rx, top, rw, Inches(0.4),
             [[("SOURCE RELIABILITY TIERS", {"size": 11.5, "bold": True, "color": GREEN})]])
    tiers = [
        ("A", "Primary / authoritative", "legislation, regulators (SAHPRA, SAPC, CMS, HPCSA), peer-reviewed, standards, company disclosures", GREEN),
        ("B", "Strong secondary", "reputable analysts & consultancies, established trade press, gov-adjacent bodies", GREEN_DK),
        ("C", "Useful but weaker", "general news, vendor marketing, single-author blogs — verify load-bearing claims", RGBColor(0xC9, 0xA2, 0x2E)),
        ("D", "Anecdotal / signal-only", "forums, app-store reviews, social posts — never a factual basis; sentiment signal only", RGBColor(0xC9, 0x6A, 0x2E)),
    ]
    yy = top + Inches(0.45)
    for letter, name, desc, col in tiers:
        add_rect(slide, rx, yy, rw, Inches(0.82), fill=WHITE, line=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        badge = add_rect(slide, rx + Inches(0.16), yy + Inches(0.16), Inches(0.5),
                         Inches(0.5), fill=col, shape=MSO_SHAPE.OVAL)
        shape_text(badge, [[(letter, {"size": 15, "bold": True, "color": WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
        add_text(slide, rx + Inches(0.8), yy + Inches(0.1), rw - Inches(1.0), Inches(0.3),
                 [[(name, {"size": 11, "bold": True, "color": INK})]])
        add_text(slide, rx + Inches(0.8), yy + Inches(0.4), rw - Inches(1.0), Inches(0.4),
                 [[(desc, {"size": 8.5, "color": SLATE})]], line_spacing=0.92)
        yy += Inches(0.88)
    notes(slide, "Appendix for the curious. The vault follows the Karpathy 'LLM wiki' "
                 "pattern in three layers: raw sources (append-only evidence), the "
                 "synthesised wiki (where claims link back to their sources), and the "
                 "schema (the rules). The golden rule is that synthesis lives in the "
                 "wiki and every claim cites its evidence. We weight that evidence on an "
                 "A-D scale: regulatory, clinical and pricing claims rest on A/B; "
                 "consumer sentiment from social and app-store data is C/D and always "
                 "labelled as a signal, never asserted as a rate.")
    return slide


def slide_appendix_personas_reform(prs):
    slide = new_slide(prs)
    content_chrome(
        slide, "Appendix", "Full persona set & open questions / reform register",
        "The complete cast the design must satisfy, and the honest list of what we "
        "still need to verify or change in regulation.",
        cite="design-kit/personas/_personas-index · meta/open-questions.md")
    top = Inches(2.4)
    # left: 8 personas list
    lx = Inches(0.6)
    add_text(slide, lx, top, Inches(6.0), Inches(0.35),
             [[("THE 8 PERSONAS", {"size": 11.5, "bold": True, "color": GREEN})]])
    personas = [
        ("Nomvula Khumalo", "Insured chronic, multimorbid ★"),
        ("Thabo Mokoena", "Cash-pay acute customer"),
        ("Lerato Dlamini", "Caregiver / proxy collector ★"),
        ("Trudie van Wyk", "Elderly poly-pharmacy patient"),
        ("Sipho Ndlovu", "Low-data, peri-urban, WhatsApp-first ★"),
        ("Aisha Patel", "Time-poor urban professional"),
        ("Mandla Zwane", "Public-sector bridging to retail"),
        ("Naledi Mahlangu", "Frontline dispensary pharmacist (staff)"),
    ]
    yy = top + Inches(0.42)
    for i, (name, role) in enumerate(personas):
        add_text(slide, lx, yy, Inches(6.0), Inches(0.38),
                 [[(f"{i+1}.  ", {"size": 10.5, "bold": True, "color": GREEN_DK}),
                   (name + "  ", {"size": 10.5, "bold": True, "color": INK}),
                   ("— " + role, {"size": 10, "color": SLATE})]])
        yy += Inches(0.42)
    add_text(slide, lx, yy + Inches(0.05), Inches(6.0), Inches(0.3),
             [[("★ = protagonist of the North Star journey",
                {"size": 8.5, "italic": True, "color": GREEN_DK})]])
    # right: open questions / reform register
    rx = Inches(6.95)
    rw = Inches(5.78)
    add_text(slide, rx, top, rw, Inches(0.35),
             [[("OPEN QUESTIONS & REFORM REGISTER (selected)",
                {"size": 11.5, "bold": True, "color": GREEN})]])
    items = [
        "Real-time scheme pre-adjudication API (“can I claim?” dry-run) — exposed by any switch? Key enabler for U3.",
        "Legal boundary of auto-dispense for eligible repeats vs pharmacist oversight. \U0001f52e",
        "SAPC acceptance of AES-signed e-scripts for S5/S6. \U0001f52e",
        "Chain-of-custody for third-party last-mile couriers holding scheduled meds. \U0001f52e",
        "24/7 pharmacist oversight & live after-hours scheme authorisation. \U0001f52e",
        "Exact dispensing-fee bands & CMS coverage split — re-verify vs primary gazette/PDF.",
    ]
    yy = top + Inches(0.42)
    for it in items:
        add_text(slide, rx, yy, rw, Inches(0.62),
                 [[("▸ ", {"size": 10, "bold": True, "color": GREEN}),
                   (it, {"size": 10, "color": INK})]], line_spacing=1.0)
        yy += Inches(0.6)
    add_text(slide, rx, yy + Inches(0.02), rw, Inches(0.35),
             [[("\U0001f52e = requires regulatory reform.  ", {"size": 8.5, "italic": True, "color": AMBER}),
               ("48 reform flags & 63 unverified items tracked across the vault.",
                {"size": 8.5, "italic": True, "color": SLATE})]], line_spacing=1.0)
    notes(slide, "Closing appendix. On the left, the full cast of eight personas the "
                 "design must satisfy — including Naledi, the frontline pharmacist, "
                 "whose oversight and time the future-state must protect, not erode. On "
                 "the right, the honesty register: the open questions and reform-"
                 "dependencies we are tracking, from real-time scheme pre-adjudication "
                 "to the legal boundary of automated dispensing and after-hours "
                 "authorisation. The vault tracks 48 reform flags and 63 unverified "
                 "items in total. Naming these openly is what makes the work "
                 "trustworthy.")
    return slide


# ============================================================================
# ASSEMBLY
# ============================================================================

def build():
    for p in (IMG_BOARD, IMG_MOMENT, IMG_CASH, IMG_CABINET):
        if not os.path.exists(p):
            raise FileNotFoundError(f"Missing required asset: {p}")

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    PAGE["n"] = 0

    # 1. Title
    slide_title(prs)
    # 2. Executive summary
    slide_exec_summary(prs)
    # 3. Why now (problem)
    slide_why_now(prs)

    # Divider: The method
    section_divider(prs, "01", "The method",
                    "How we built a defensible North Star — and the evidence it stands on.")
    # 4. Approach / pipeline
    slide_approach(prs)
    # 5. Evidence base
    slide_evidence(prs)
    # 6. Best practice
    slide_best_practice(prs)

    # Divider: Principles & people
    section_divider(prs, "02", "Principles & people",
                    "The strategic bets, the experience constitution, the segments and the personas the journey is grounded in.")
    # 7. Strategic principles
    slide_strategic(prs)
    # 8. UX principles
    slide_ux(prs)
    # 9. Segments
    slide_segments(prs)
    # 10. Personas
    slide_personas(prs)
    # 11. Pains mapped
    slide_pains_mapped(prs)

    # Divider: The North Star
    section_divider(prs, "03", "The North Star",
                    "The future-state experience — the household medicine manager, revealed.")
    # 12-14. North Star reveal (3 slides)
    slide_northstar_intro(prs)
    slide_northstar_moment(prs)
    slide_northstar_cabinet(prs)
    # 15. What makes it real
    slide_makes_it_real(prs)
    # 16. From pain to North Star
    slide_before_after(prs)

    # Divider: From here
    section_divider(prs, "04", "From here",
                    "What this deck is, what comes next, and the evidence behind it.")
    # 17. What's next
    slide_whats_next(prs)
    # 18. Appendix — method & tiers
    slide_appendix_method(prs)
    # 19. Appendix — personas & reform
    slide_appendix_personas_reform(prs)

    prs.save(OUT)
    return OUT, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    path, count = build()
    size = os.path.getsize(path)
    print(f"Wrote {path}")
    print(f"Size: {size:,} bytes ({size/1024:.0f} KB)")
    print(f"Slides: {count}")
