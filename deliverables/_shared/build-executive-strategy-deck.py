#!/usr/bin/env python3
"""
build-executive-strategy-deck.py
================================
Re-runnable, offline generator for the **Executive Scripting Strategy Deck**

    deliverables/00-scripting-strategy-deck.pptx   (+ .pdf via LibreOffice)

"Future-State Scripting — A Strategy for Dis-Chem"  (Bigly Labs x Dis-Chem)

This engine EXTENDS the North Star deck pipeline (build-northstar-deck.py): it
copies that deck's Dis-Chem palette and core drawing helpers verbatim, then adds
a catalogue of *parameterised* layout renderers (SPEC §2) that consume content
authored in Markdown (SPEC §3). Content agents write
`deliverables/_shared/exec-deck/content/part{0..6}-*.md`; this engine renders
them in part/slide order into the deck.

It is DEFENSIVE by design: it never crashes on a missing/empty/partially-written
content file, an unknown layout, or a missing field. Unknowns are rendered as
clearly-labelled warning slides and echoed to stderr; if `content/` is empty it
still produces a cover slide + a "no content yet" notice so the pipeline runs.

Dependencies:  python-pptx 1.0.2  (pip install python-pptx). No pyyaml needed —
the content loader is a small hand-rolled tolerant parser. No network.

Run:
    python3 deliverables/_shared/build-executive-strategy-deck.py

PDF step (run after the .pptx is built):
    libreoffice --headless --convert-to pdf deliverables/00-scripting-strategy-deck.pptx --outdir deliverables/
"""

import os
import re
import sys
import glob

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Paths
# ----------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "deck-assets")
EXEC_DECK = os.path.join(HERE, "exec-deck")
CONTENT_DIR = os.path.join(EXEC_DECK, "content")
OUT = os.path.abspath(os.path.join(HERE, "..", "00-scripting-strategy-deck.pptx"))

# Visual id -> embedded PNG file (SPEC §4 "img:*" catalogue)
IMG_MAP = {
    "medicine-cabinet": os.path.join(ASSETS, "medicine-cabinet.png"),
    "journey-board": os.path.join(ASSETS, "journey-board.png"),
    "journey-cash": os.path.join(ASSETS, "journey-cash.png"),
    "journey-moment-expanded": os.path.join(ASSETS, "journey-moment-expanded.png"),
}

# ----------------------------------------------------------------------------
# Palette (Dis-Chem) — copied verbatim from build-northstar-deck.py
# ----------------------------------------------------------------------------
GREEN = RGBColor(0x1A, 0x9B, 0x4A)   # primary accent / dividers
GREEN_DK = RGBColor(0x12, 0x6E, 0x35)
LIME = RGBColor(0xB5, 0xD3, 0x34)    # sparing accent
INK = RGBColor(0x1D, 0x1D, 0x1D)     # headings / dark text
SLATE = RGBColor(0x52, 0x52, 0x5B)   # body text
MINT = RGBColor(0xE8, 0xF5, 0xEC)    # light fill
MINT_MD = RGBColor(0xBF, 0xE3, 0xCF)  # mid mint fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xFA, 0xFB, 0xFA)   # near-white bg
LINE = RGBColor(0xD7, 0xDD, 0xD9)    # hairline
AMBER = RGBColor(0xB0, 0x6A, 0x00)   # reform / caution accent
CARD_BG = RGBColor(0xF4, 0xF8, 0xF5)
RED_DK = RGBColor(0x9B, 0x37, 0x37)  # "today / pain" accent

FONT = "Calibri"
FONT_LT = "Calibri Light"

# Slide geometry (16:9)
SW = Inches(13.333)
SH = Inches(7.5)

FOOTER = ("Concept for review — not an official Dis-Chem artefact  ·  "
          "Bigly Labs × Dis-Chem")

# ----------------------------------------------------------------------------
# Low-level helpers — copied verbatim from build-northstar-deck.py
# ----------------------------------------------------------------------------

def _set_bg(slide, color):
    """Solid slide background via a full-bleed rectangle sent to back."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    sp = rect._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return rect


def _no_autofit(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0):
    """runs: list of paragraphs; each paragraph is a list of (text, opts) tuples.
    opts: dict with size, bold, italic, color, font, underline."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        for text, opts in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", FONT)
            f.size = Pt(opts.get("size", 14))
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
            f.underline = opts.get("underline", False)
            f.color.rgb = opts.get("color", INK)
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
             shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    return sp


def _soft_shadow(shape):
    """Add a subtle outer shadow to a shape."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    shdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': '60000', 'dist': '25000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = shdw.makeelement(qn('a:srgbClr'), {'val': '1D1D1D'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '22000'})
    clr.append(alpha)
    shdw.append(clr)
    effectLst.append(shdw)
    spPr.append(effectLst)


def shape_text(shape, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
               margins=0.06):
    tf = shape.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margins)
    tf.margin_right = Inches(margins)
    tf.margin_top = Inches(margins * 0.6)
    tf.margin_bottom = Inches(margins * 0.6)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.0
        for text, opts in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = opts.get("font", FONT)
            f.size = Pt(opts.get("size", 12))
            f.bold = opts.get("bold", False)
            f.italic = opts.get("italic", False)
            f.color.rgb = opts.get("color", INK)
    return shape


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ----------------------------------------------------------------------------
# Diagram geometry + drawing helpers (WS4)
# ----------------------------------------------------------------------------
# Native-diagram content canvas (diagrams.md "Shared conventions"). A native
# diagram OWNS the area below the title/so-what band. We start a touch below the
# spec's y=1.7 so the diagram never collides with the engine's so-what strip
# (which sits at y≈1.62→2.14). Everything a diagram draws must stay inside this
# box; helpers below clamp/assert against it.
DIA_X0 = 0.6
DIA_X1 = 12.73
DIA_Y0 = 2.32      # below the so-what band
DIA_Y1 = 6.05      # above the flag-strip / footer
DIA_W = DIA_X1 - DIA_X0
DIA_H = DIA_Y1 - DIA_Y0

# Hard slide bounds (inches) for the defensive clamp/assert.
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Live diagram canvas transform. Diagram functions are written against the
# *spec* canvas (DIA_X0..DIA_X1 / DIA_Y0..DIA_Y1). When a diagram must render
# into a smaller region (e.g. beside a fact column), we set this transform so
# every helper maps spec-inches -> actual-inches and scales font sizes by `s`.
# Default = identity (full canvas). Use the `diagram_region()` context manager.
_CANVAS = {"ox": 0.0, "oy": 0.0, "sx": 1.0, "sy": 1.0, "s": 1.0,
           "ax0": DIA_X0, "ay0": DIA_Y0, "ax1": DIA_X1, "ay1": DIA_Y1}


class diagram_region:
    """Context manager: remap the spec diagram canvas into an actual sub-box
    (inches). Scales positions and font sizes uniformly (min of sx,sy)."""

    def __init__(self, x, y, w, h):
        self.box = (x, y, w, h)

    def __enter__(self):
        x, y, w, h = self.box
        sx = w / DIA_W
        sy = h / DIA_H
        s = min(sx, sy)
        _CANVAS.update(ox=x - DIA_X0 * sx, oy=y - DIA_Y0 * sy,
                       sx=sx, sy=sy, s=s,
                       ax0=x, ay0=y, ax1=x + w, ay1=y + h)
        return self

    def __exit__(self, *exc):
        _CANVAS.update(ox=0.0, oy=0.0, sx=1.0, sy=1.0, s=1.0,
                       ax0=DIA_X0, ay0=DIA_Y0, ax1=DIA_X1, ay1=DIA_Y1)
        return False


def _tx(x):
    return _CANVAS["ox"] + x * _CANVAS["sx"]


def _ty(y):
    return _CANVAS["oy"] + y * _CANVAS["sy"]


def _tw(w):
    return w * _CANVAS["sx"]


def _th(h):
    return h * _CANVAS["sy"]


# Minimum on-slide font size for any diagram label (projector legibility floor).
# When a diagram is drawn into a scaled sub-region the raw multiply can drop
# labels far below this; _ts() floors the *scaled* size so no diagram text ever
# renders below DIA_FONT_FLOOR pt. At full canvas (s==1.0) this is a no-op for
# any authored size >= the floor.
DIA_FONT_FLOOR = 9.0


def _ts(size):
    scaled = size * _CANVAS["s"]
    if _CANVAS["s"] != 1.0 and scaled < DIA_FONT_FLOOR:
        return DIA_FONT_FLOOR
    return scaled


def _scale_runs(runs):
    """Scale font sizes inside a runs structure by the active canvas scale."""
    if _CANVAS["s"] == 1.0:
        return runs
    out = []
    for para in runs:
        np = []
        for text, opts in para:
            o = dict(opts)
            o["size"] = _ts(o.get("size", 12))
            np.append((text, o))
        out.append(np)
    return out


def _clamp(v, lo, hi):
    return lo if v < lo else (hi if v > hi else v)


def _fit_box(x, y, w, h, x0=None, y0=None, x1=None, y1=None):
    """Clamp a box (spec inches) to the active canvas; never returns negative
    width/height. Defensive: keeps every diagram shape inside the canvas even if
    a spec coordinate is slightly off."""
    x0 = DIA_X0 if x0 is None else x0
    y0 = DIA_Y0 if y0 is None else y0
    x1 = DIA_X1 if x1 is None else x1
    y1 = DIA_Y1 if y1 is None else y1
    x = _clamp(x, x0, x1)
    y = _clamp(y, y0, y1)
    if x + w > x1:
        w = max(0.05, x1 - x)
    if y + h > y1:
        h = max(0.05, y1 - y)
    return x, y, w, h


def dia_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=False, clamp=True):
    """add_rect in SPEC inches, transformed to actual canvas + clamped."""
    if clamp:
        x, y, w, h = _fit_box(x, y, w, h)
    return add_rect(slide, Inches(_tx(x)), Inches(_ty(y)), Inches(_tw(w)),
                    Inches(_th(h)), fill=fill, line=line, line_w=line_w,
                    shape=shape, shadow=shadow)


def dia_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, clamp=True):
    """add_text in SPEC inches, transformed + clamped; font sizes scaled."""
    if clamp:
        x, y, w, h = _fit_box(x, y, w, h, y1=DIA_Y1 + 0.25)
    return add_text(slide, Inches(_tx(x)), Inches(_ty(y)), Inches(_tw(w)),
                    Inches(_th(h)), _scale_runs(runs),
                    align=align, anchor=anchor, line_spacing=line_spacing)


def chip(slide, x, y, w, h, label, sub=None, fill=MINT, txt=GREEN_DK,
         line=None, line_w=None, size=10.5, sub_size=8.5, bold=True,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER, clamp=True):
    """A labelled rounded chip/node with optional sub-caption. SPEC inches."""
    if clamp:
        x, y, w, h = _fit_box(x, y, w, h)
    box = add_rect(slide, Inches(_tx(x)), Inches(_ty(y)), Inches(_tw(w)),
                   Inches(_th(h)), fill=fill, line=line, line_w=line_w, shape=shape)
    runs = [[(label, {"size": _ts(size), "bold": bold, "color": txt})]]
    if sub:
        runs.append([(sub, {"size": _ts(sub_size), "bold": False, "color": txt})])
    shape_text(box, runs, align=align, anchor=MSO_ANCHOR.MIDDLE,
               margins=0.07 * _CANVAS["s"])
    return box


def _line_xml(shape, color, width_pt, dashed=False, head=False, tail=False):
    """Style a connector line element: colour, width, optional dash + arrowheads."""
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)
    lnEl = shape._element.spPr.find(qn('a:ln'))
    if lnEl is None:
        lnEl = shape._element.spPr.makeelement(qn('a:ln'), {})
        shape._element.spPr.append(lnEl)
    # remove pre-existing dash / head / tail
    for tag in ('a:prstDash', 'a:headEnd', 'a:tailEnd'):
        for el in lnEl.findall(qn(tag)):
            lnEl.remove(el)
    if dashed:
        d = lnEl.makeelement(qn('a:prstDash'), {'val': 'dash'})
        lnEl.append(d)
    if head:
        h = lnEl.makeelement(qn('a:headEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'})
        lnEl.append(h)
    if tail:
        t = lnEl.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'})
        lnEl.append(t)


def connector(slide, x1, y1, x2, y2, color=GREEN, width=1.5,
              dashed=False, arrow=True, two_way=False, clamp=True):
    """Straight connector (SPEC inches), transformed; tail arrowhead (head too
    if two_way). Line width scales with the canvas."""
    if clamp:
        x1 = _clamp(x1, DIA_X0, DIA_X1); x2 = _clamp(x2, DIA_X0, DIA_X1)
        y1 = _clamp(y1, DIA_Y0, DIA_Y1); y2 = _clamp(y2, DIA_Y0, DIA_Y1)
    cxn = slide.shapes.add_connector(
        2, Inches(_tx(x1)), Inches(_ty(y1)), Inches(_tx(x2)), Inches(_ty(y2)))
    cxn.shadow.inherit = False
    _line_xml(cxn, color, max(0.5, width * _CANVAS["s"]), dashed=dashed,
              tail=arrow, head=(arrow and two_way))
    return cxn


def chevron(slide, x, y, w, h, fill=GREEN, label=None, txt=WHITE, size=11,
            clamp=True):
    """A right-pointing chevron block (SPEC inches), optional centred label."""
    if clamp:
        x, y, w, h = _fit_box(x, y, w, h)
    sp = add_rect(slide, Inches(_tx(x)), Inches(_ty(y)), Inches(_tw(w)),
                  Inches(_th(h)), fill=fill, shape=MSO_SHAPE.CHEVRON)
    if label:
        shape_text(sp, [[(label, {"size": _ts(size), "bold": True, "color": txt})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.04)
    return sp


def dia_caption(slide, text, y=None, color=SLATE, size=9, italic=True,
                align=PP_ALIGN.LEFT):
    """A full-width diagram caption line near the canvas bottom."""
    if text is None:
        return
    if y is None:
        y = DIA_Y1 - 0.02
    dia_text(slide, DIA_X0, y, DIA_W, 0.4,
             [[(text, {"size": size, "italic": italic, "color": color})]],
             align=align, clamp=False)


def person_glyph(slide, cx, cy, scale=0.45, color=SLATE, label=None,
                 label_color=None):
    """Simple person icon centred at (cx, cy) in SPEC inches (head + body)."""
    head_d = 0.38 * scale / 0.45
    body_w = 0.62 * scale / 0.45
    body_h = 0.52 * scale / 0.45
    add_rect(slide, Inches(_tx(cx - head_d / 2)),
             Inches(_ty(cy - body_h / 2 - head_d * 0.9)),
             Inches(_tw(head_d)), Inches(_th(head_d)), fill=color,
             shape=MSO_SHAPE.OVAL)
    add_rect(slide, Inches(_tx(cx - body_w / 2)),
             Inches(_ty(cy - body_h / 2 + head_d * 0.1)),
             Inches(_tw(body_w)), Inches(_th(body_h)), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if label:
        dia_text(slide, cx - 0.8, cy + body_h / 2 + head_d * 0.15, 1.6, 0.26,
                 [[(label, {"size": 9, "bold": True,
                            "color": label_color or color})]],
                 align=PP_ALIGN.CENTER)


def tick_glyph(slide, cx, cy, kind="check", d=0.30):
    """Centre a tick/cross/dash glyph at (cx,cy) SPEC inches. check|cross|dash."""
    if kind == "check":
        circ = add_rect(slide, Inches(_tx(cx - d / 2)), Inches(_ty(cy - d / 2)),
                        Inches(_tw(d)), Inches(_th(d)), fill=GREEN,
                        shape=MSO_SHAPE.OVAL)
        shape_text(circ, [[("✓", {"size": _ts(11), "bold": True, "color": WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.0)
    elif kind == "cross":
        dia_text(slide, cx - d / 2, cy - d / 2 - 0.04, d, d + 0.1,
                 [[("✗", {"size": 13, "bold": True, "color": SLATE})]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:  # dash / partial
        dia_text(slide, cx - d / 2, cy - d / 2 - 0.04, d, d + 0.1,
                 [[("–", {"size": 14, "bold": True, "color": SLATE})]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def draw_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
              shape=MSO_SHAPE.RECTANGLE, shadow=False):
    """Transformed raw rectangle in SPEC inches (no clamping; for fills/strips
    that diagrams position precisely). Returns the shape."""
    return add_rect(slide, Inches(_tx(x)), Inches(_ty(y)), Inches(_tw(w)),
                    Inches(_th(h)), fill=fill, line=line, line_w=line_w,
                    shape=shape, shadow=shadow)


def draw_label(shape, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
               margins=0.07):
    """shape_text with font sizes scaled to the active canvas."""
    return shape_text(shape, _scale_runs(runs), align=align, anchor=anchor,
                      margins=margins * _CANVAS["s"])


# ----------------------------------------------------------------------------
# Slide chrome (content slides) — copied verbatim from build-northstar-deck.py
# ----------------------------------------------------------------------------
PAGE = {"n": 0}


def content_chrome(slide, kicker, title, takeaway, cite=None):
    """Standard content-slide header: kicker, title, accent rule, so-what strip."""
    _set_bg(slide, WHITE)
    add_rect(slide, 0, 0, SW, Inches(0.12), fill=GREEN)
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.34), Inches(11.5), Inches(0.3),
                 [[(kicker.upper(), {"size": 11, "bold": True, "color": GREEN,
                                     "font": FONT})]])
    add_text(slide, Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.9),
             [[(title, {"size": 27, "bold": True, "color": INK})]],
             line_spacing=0.98)
    if takeaway:
        sy = Inches(1.62)
        add_rect(slide, Inches(0.6), sy, Inches(0.09), Inches(0.52), fill=GREEN)
        band = add_rect(slide, Inches(0.69), sy, Inches(12.04), Inches(0.52), fill=MINT)
        shape_text(band, [[("So what:  ", {"size": 12.5, "bold": True, "color": GREEN_DK}),
                           (takeaway, {"size": 12.5, "bold": False, "color": INK})]],
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margins=0.16)
    if cite:
        add_text(slide, Inches(0.6), Inches(7.06), Inches(8.0), Inches(0.3),
                 [[("Source: ", {"size": 8, "italic": True, "color": SLATE}),
                   (cite, {"size": 8, "italic": True, "color": GREEN_DK})]])
    _footer(slide)


def _footer(slide):
    PAGE["n"] += 1
    add_rect(slide, Inches(0.0), Inches(7.18), SW, Pt(0.75), fill=LINE)
    add_text(slide, Inches(0.6), Inches(7.24), Inches(11.0), Inches(0.24),
             [[(FOOTER, {"size": 7.5, "color": SLATE})]])
    add_text(slide, Inches(12.4), Inches(7.24), Inches(0.7), Inches(0.24),
             [[(str(PAGE["n"]), {"size": 8, "bold": True, "color": GREEN_DK})]],
             align=PP_ALIGN.RIGHT)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ----------------------------------------------------------------------------
# Section divider — copied verbatim from build-northstar-deck.py
# ----------------------------------------------------------------------------

def section_divider(prs, num, title, subtitle):
    slide = new_slide(prs)
    _set_bg(slide, GREEN)
    add_rect(slide, Inches(11.7), Inches(0.7), Inches(0.5), Inches(2.4), fill=GREEN_DK)
    add_rect(slide, Inches(10.75), Inches(1.65), Inches(2.4), Inches(0.5), fill=GREEN_DK)
    add_text(slide, Inches(0.9), Inches(2.7), Inches(2.0), Inches(0.6),
             [[(num, {"size": 40, "bold": True, "color": LIME, "font": FONT_LT})]])
    add_rect(slide, Inches(0.95), Inches(3.55), Inches(0.9), Inches(0.06), fill=LIME)
    add_text(slide, Inches(0.9), Inches(3.75), Inches(11.0), Inches(1.4),
             [[(title, {"size": 34, "bold": True, "color": WHITE})]], line_spacing=1.0)
    if subtitle:
        add_text(slide, Inches(0.95), Inches(5.0), Inches(10.6), Inches(1.0),
                 [[(subtitle, {"size": 14, "color": MINT_MD})]], line_spacing=1.15)
    PAGE["n"] += 1
    add_text(slide, Inches(12.4), Inches(7.24), Inches(0.7), Inches(0.24),
             [[(str(PAGE["n"]), {"size": 8, "bold": True, "color": MINT_MD})]],
             align=PP_ALIGN.RIGHT)
    return slide


# ----------------------------------------------------------------------------
# Shared building blocks (reused across layouts)
# ----------------------------------------------------------------------------

def _bullet_block(slide, x, y, w, items, size=12.5, color=INK, gap=0.82,
                  marker="▸ ", marker_color=GREEN):
    yy = y
    for it in items:
        add_text(slide, x, yy, w, Inches(0.78),
                 [[(marker, {"size": size, "bold": True, "color": marker_color}),
                   (it, {"size": size, "color": color})]], line_spacing=1.05)
        yy += Inches(gap)


def _pillar_grid(slide, items, top, cols, accent, card_h=None, small=False):
    """items: list of (id, name, desc) tuples."""
    n = len(items)
    if n == 0:
        return
    gx = Inches(0.22)
    gy = Inches(0.22)
    lx = Inches(0.6)
    avail_w = SW - lx * 2
    cw = (avail_w - gx * (cols - 1)) / cols
    if card_h is None:
        card_h = Inches(1.72)
    for idx, (pid, name, desc) in enumerate(items):
        r, c = divmod(idx, cols)
        x = lx + c * (cw + gx)
        y = top + r * (card_h + gy)
        add_rect(slide, x, y, cw, card_h, fill=CARD_BG, line=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        if pid:
            badge = add_rect(slide, x + Inches(0.18), y + Inches(0.18), Inches(0.62),
                             Inches(0.42), fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            shape_text(badge, [[(pid, {"size": 14, "bold": True, "color": WHITE})]],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
            name_x = x + Inches(0.95)
            name_w = cw - Inches(1.1)
        else:
            name_x = x + Inches(0.2)
            name_w = cw - Inches(0.4)
        nsize = 11.5 if not small else 11
        dsize = 9.5 if not small else 9
        add_text(slide, name_x, y + Inches(0.16), name_w, Inches(0.7),
                 [[(name, {"size": nsize, "bold": True, "color": INK})]],
                 line_spacing=0.95)
        if desc:
            add_text(slide, x + Inches(0.2), y + Inches(0.78 if not small else 0.74),
                     cw - Inches(0.4), card_h - Inches(0.9),
                     [[(desc, {"size": dsize, "color": SLATE})]], line_spacing=1.0)


def _flag_strip(slide, flags, y=Inches(6.78)):
    """Render honesty flags (unverified / reform) as a small on-slide footnote."""
    if not flags:
        return
    parts = []
    for fl in flags:
        f = str(fl).strip().lower()
        if f == "unverified":
            parts.append(("⚠ vendor-reported / extract-only (unverified)", AMBER))
        elif f == "reform":
            parts.append(("\U0001f52e requires regulatory reform", AMBER))
        elif f:
            parts.append((f, SLATE))
    if not parts:
        return
    runs = []
    for i, (txt, col) in enumerate(parts):
        if i:
            runs.append(("    ", {"size": 9, "color": SLATE}))
        runs.append((txt, {"size": 9, "italic": True, "color": col}))
    add_text(slide, Inches(0.6), y, Inches(11.5), Inches(0.3), [runs])


def _evidence_line(items):
    """Format an evidence[] list (vault paths/slugs) into a single cite string."""
    if not items:
        return None
    if isinstance(items, str):
        return items
    return " · ".join(str(i) for i in items)


# ============================================================================
# NATIVE DIAGRAMS (WS4) — each draws into the full diagram canvas (SPEC §4 /
# diagrams.md). Signature: fn(slide, caption=None). The slide already has its
# title + so-what band drawn by content_chrome; the diagram owns y≈2.32→6.05.
# Conventions (diagrams.md): solid=operable · dashed-AMBER=reform · GREEN=desired
# path · SLATE/LINE/PAPER=legacy · max ONE LIME element per diagram.
# ============================================================================

def dia_front_door(slide, caption=None):
    """Counter → front door (left muted TODAY, right accented FRONT DOOR)."""
    # LEFT — TODAY
    dia_rect(slide, 0.6, 2.45, 5.0, 3.05, fill=PAPER, line=LINE, line_w=Pt(1),
             shape=MSO_SHAPE.RECTANGLE)
    chip(slide, 0.78, 2.55, 1.1, 0.32, "TODAY", fill=None, txt=SLATE, size=10,
         shape=MSO_SHAPE.RECTANGLE)
    chip(slide, 1.7, 2.95, 2.8, 0.78, "Dispensing counter", sub="one-off transaction",
         fill=WHITE, txt=INK, line=LINE, line_w=Pt(1), size=11, sub_size=9)
    for i, lab in enumerate(("Queue", "Hand over script", "Pay & leave")):
        chip(slide, 1.5, 3.92 + i * 0.46, 3.2, 0.38, lab, fill=CARD_BG, txt=SLATE,
             line=LINE, line_w=Pt(0.75), size=9.5, bold=False)
    # CHEVRON
    chevron(slide, 5.78, 3.55, 0.95, 0.7, fill=GREEN)
    # RIGHT — FRONT DOOR
    dia_rect(slide, 6.9, 2.45, 5.83, 3.05, fill=MINT, line=GREEN, line_w=Pt(1.25),
             shape=MSO_SHAPE.RECTANGLE)
    chip(slide, 7.06, 2.55, 1.9, 0.32, "THE FRONT DOOR", fill=None, txt=GREEN_DK,
         size=10, shape=MSO_SHAPE.RECTANGLE)
    hub = chip(slide, 9.0, 2.96, 2.5, 0.62, "Ongoing care relationship",
               fill=GREEN, txt=WHITE, size=11.5)
    sats = ["Screening", "Advice", "Chronic management", "Telehealth",
            "Behaviour change"]
    sat_y = 3.92
    sat_x = [7.05, 8.42, 9.79, 11.16, 7.74]
    sat_w = 1.30
    # row 1 (4 chips) + row 2 (1 centred)
    positions = [(7.05, sat_y), (8.42, sat_y), (9.79, sat_y), (11.16, sat_y),
                 (9.0, sat_y + 0.62)]
    for (sx, sy), lab in zip(positions, sats):
        connector(slide, 10.25, 3.58, sx + sat_w / 2, sy + 0.18, color=GREEN,
                  width=1.0, arrow=False)
        chip(slide, sx, sy, sat_w, 0.46, lab, fill=MINT_MD, txt=GREEN_DK, size=9)
    dia_caption(slide, caption or ("~40% of consumers want broader clinical services "
                "via the pharmacy — direction of travel."), y=5.55)


def dia_fulfilment_node(slide, caption=None):
    """Hub-and-spoke channel ladder: inbound channels → hub → outbound modes."""
    hub = chip(slide, 5.0, 2.85, 3.3, 1.4, "The store as fulfilment node",
               sub="one inventory · one record", fill=GREEN, txt=WHITE, size=12,
               sub_size=9)
    inbound = ["WhatsApp / USSD / SMS", "App", "Walk-in / counter"]
    outbound = [("Collect in store", False), ("Home delivery", False),
                ("Locker / pickup point", True)]
    ys = [2.55, 3.45, 4.35]
    for lab, y in zip(inbound, ys):
        chip(slide, 0.7, y, 2.6, 0.62, lab, fill=MINT, txt=GREEN_DK, size=10)
        connector(slide, 3.3, y + 0.31, 5.0, 3.55, color=GREEN, width=1.25)
    dia_text(slide, 0.7, 5.02, 4.3, 0.34,
             [[("channel ladder — meet them where they are",
                {"size": 9, "italic": True, "color": SLATE})]])
    for (lab, reform), y in zip(outbound, ys):
        if reform:
            box = chip(slide, 9.4, y, 2.85, 0.62, "🔮 " + lab, fill=MINT_MD,
                       txt=GREEN_DK, size=9.5, line=AMBER, line_w=Pt(1.5))
            _line_xml(box, AMBER, 1.5, dashed=True)
        else:
            chip(slide, 9.4, y, 2.85, 0.62, lab, fill=MINT_MD, txt=GREEN_DK, size=10)
        connector(slide, 8.3, 3.55, 9.4, y + 0.31, color=GREEN, width=1.25)
    dia_text(slide, 9.4, 5.02, 3.3, 0.34,
             [[("scheduled-med lockers need reform",
                {"size": 9, "italic": True, "color": AMBER})]])
    band = dia_rect(slide, 0.6, 5.42, 12.13, 0.52, fill=CARD_BG, line=LINE,
                    line_w=Pt(0.75), shape=MSO_SHAPE.RECTANGLE)
    draw_label(band, [[(caption or ("The convenience bar is set by retail (Sixty60), "
                        "not other pharmacies · only ~1 in 10 UK patients manage "
                        "repeats digitally — the adoption gap."),
                        {"size": 9, "color": INK})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.12)


def dia_time_reallocation(slide, caption=None):
    """Two stacked time-budget bars: toil shrinks (SLATE→LINE), clinical grows
    (MINT_MD→GREEN)."""
    bx, bw = 1.6, 10.0
    # Bar 1 — TODAY
    dia_text(slide, 0.6, 2.5, 0.95, 0.4,
             [[("TODAY", {"size": 10, "bold": True, "color": SLATE})]],
             anchor=MSO_ANCHOR.MIDDLE)
    s1 = bw * 0.70
    seg = draw_rect(slide, bx, 2.45, s1, 0.85, fill=SLATE)
    draw_label(seg, [[("Counting · checking · admin (toil)",
                       {"size": 10, "bold": True, "color": WHITE})]])
    seg = draw_rect(slide, bx + s1, 2.45, bw - s1, 0.85, fill=MINT_MD)
    draw_label(seg, [[("Clinical · counselling", {"size": 9.5, "color": INK})]])
    # transition
    chevron(slide, 6.3, 3.42, 0.5, 0.42, fill=GREEN)
    dia_text(slide, 6.9, 3.46, 5.0, 0.34,
             [[("automation + AI absorb the toil",
                {"size": 9, "bold": True, "color": GREEN_DK})]],
             anchor=MSO_ANCHOR.MIDDLE)
    # Bar 2 — FUTURE
    dia_text(slide, 0.6, 4.05, 0.95, 0.4,
             [[("FUTURE", {"size": 10, "bold": True, "color": GREEN_DK})]],
             anchor=MSO_ANCHOR.MIDDLE)
    s2 = bw * 0.25
    seg = draw_rect(slide, bx, 4.0, s2, 0.85, fill=LINE)
    draw_label(seg, [[("Automated dispense", {"size": 9, "color": SLATE})]])
    seg = draw_rect(slide, bx + s2, 4.0, bw - s2, 0.85, fill=GREEN)
    draw_label(seg, [[("Manage chronic disease · counsel · top-of-licence",
                       {"size": 10, "bold": True, "color": WHITE})]])
    dia_caption(slide, caption or ("~70% of Danish community pharmacies already run "
                "dispensing robots · but 81% of UK pharmacists considered leaving in "
                "2022 — the workforce is the live risk."), y=5.2)


def dia_efficacy_effectiveness(slide, caption=None):
    """Two-bar gap chart; LIME bracket names the gap = adherence."""
    base_y = 5.2  # bar baseline
    # Bar A — Efficacy (full height)
    a_h = 3.0
    draw_rect(slide, 2.2, base_y - a_h, 1.8, a_h, fill=GREEN)
    dia_text(slide, 1.9, base_y - a_h - 0.5, 2.4, 0.46,
             [[("Efficacy", {"size": 12, "bold": True, "color": GREEN_DK})],
              [("what the drug can do", {"size": 9, "color": SLATE})]],
             align=PP_ALIGN.CENTER)
    # Bar B — Effectiveness (shorter)
    b_h = 2.2
    draw_rect(slide, 5.2, base_y - b_h, 1.8, b_h, fill=MINT_MD)
    dia_text(slide, 4.9, base_y - b_h - 0.5, 2.4, 0.46,
             [[("Effectiveness", {"size": 12, "bold": True, "color": INK})],
              [("what it does in real life", {"size": 9, "color": SLATE})]],
             align=PP_ALIGN.CENTER)
    # LIME gap bracket (the single highlight) — spans Bar B's shortfall
    gx = 4.6
    draw_rect(slide, gx, base_y - a_h, 0.05, a_h - b_h, fill=LIME)
    draw_rect(slide, gx, base_y - a_h, 0.4, 0.045, fill=LIME)
    draw_rect(slide, gx, base_y - b_h - 0.045, 0.4, 0.045, fill=LIME)
    # callout box
    cb = dia_rect(slide, 7.4, 2.6, 4.9, 1.4, fill=MINT, line=GREEN, line_w=Pt(1.25))
    draw_label(cb, [[("This gap = adherence", {"size": 13, "bold": True,
                                               "color": GREEN_DK})],
                    [("the experience after dispensing — the artefact in the "
                      "patient's hand", {"size": 9.5, "color": INK})]],
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margins=0.16)
    dia_text(slide, 0.6, 5.3, 12.13, 0.35,
             [[("Adherence packaging: 71% vs 63% (blister strongest, meta-analytic) "
                "— works best pharmacy-delivered & contact-rich.",
                {"size": 9, "italic": True, "color": SLATE})]])
    dia_text(slide, 0.6, 5.62, 12.13, 0.3,
             [[("⚠ vendor app/lift figures unverified",
                {"size": 9, "italic": True, "color": AMBER})]])


def dia_script_artefact(slide, caption=None):
    """Stylised redesigned-label card with three callout tags."""
    cx, cy, cw, ch = 3.85, 2.45, 5.6, 3.0
    dia_rect(slide, cx, cy, cw, ch, fill=WHITE, line=GREEN_DK, line_w=Pt(1.5))
    # header
    hdr = draw_rect(slide, cx + 0.06, cy + 0.06, cw - 0.12, 0.48, fill=GREEN)
    draw_label(hdr, [[("YOUR DAY  ", {"size": 12, "bold": True, "color": WHITE}),
                      ("⊕ Dis-Chem", {"size": 9, "color": MINT})]],
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margins=0.12)
    rows = ["☀  Morning — 2 tablets", "⛅  Midday — 1 tablet",
            "🌆  Evening — 1 tablet", "🌙  Night — 1 tablet"]
    ry = cy + 0.6
    rh = 0.44
    for i, r in enumerate(rows):
        bg = CARD_BG if i % 2 else WHITE
        cell = draw_rect(slide, cx + 0.06, ry, cw - 0.12, rh, fill=bg, line=LINE,
                         line_w=Pt(0.5))
        draw_label(cell, [[(r, {"size": 11, "bold": True, "color": INK})]],
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margins=0.18)
        ry += rh
    ftr = draw_rect(slide, cx + 0.06, ry + 0.02, cw - 0.12, 0.38, fill=MINT)
    draw_label(ftr, [[("One picture of the whole regimen",
                       {"size": 9, "bold": True, "color": GREEN_DK})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # callout tags + leaders
    t1 = chip(slide, 0.7, 2.55, 2.7, 0.6, "Icon-led — works at low literacy",
              fill=MINT_MD, txt=GREEN_DK, size=9.5)
    connector(slide, 3.4, 2.85, cx, cy + 0.85, color=GREEN, width=1.0, arrow=False)
    t2 = chip(slide, 0.7, 3.95, 2.7, 0.6, "Multilingual by default", fill=MINT_MD,
              txt=GREEN_DK, size=9.5)
    connector(slide, 3.4, 4.25, cx, cy + 1.7, color=GREEN, width=1.0, arrow=False)
    t3 = chip(slide, 9.75, 2.85, 2.95, 0.7,
              "Same artefact on paper · SMS · WhatsApp · app", fill=MINT_MD,
              txt=GREEN_DK, size=9.5)
    connector(slide, 9.75, 3.2, cx + cw, cy + 0.85, color=GREEN, width=1.0,
              arrow=False)
    dia_caption(slide, caption or ("Pattern: CVS ScriptPath (Deborah Adler, creator "
                "of Target ClearRx) — consolidates timing & dosing for polypharmacy."),
                y=5.6)


def dia_transfer_fit(slide, caption=None):
    """Ticked capability matrix — 4 condition cols × 6 region rows."""
    x0, y0 = 0.7, 2.4
    label_w = 2.4
    cols = ["Mixed funding", "Cash patients", "Store as asset",
            "Pragmatic verification"]
    col_sub = ["no single payer", "cash-pay is normal", "dense estate = node",
               "no national e-Rx rail"]
    cw = (12.6 - (x0 + label_w)) / len(cols)
    hh = 0.62
    # header
    draw_rect(slide, x0, y0, label_w, hh, fill=GREEN_DK, line=WHITE, line_w=Pt(1))
    for c, (lab, sub) in enumerate(zip(cols, col_sub)):
        cell = draw_rect(slide, x0 + label_w + c * cw, y0, cw, hh, fill=GREEN_DK,
                         line=WHITE, line_w=Pt(1))
        draw_label(cell, [[(lab, {"size": 9.5, "bold": True, "color": WHITE})],
                          [(sub, {"size": 9, "color": MINT})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.05)
    # rows: (region, strength, fill, [glyphs])  glyph: c=check x=cross d=dash
    rows = [
        ("Australia", "Highest", MINT, ["c", "c", "c", "c"]),
        ("India", "High", MINT, ["c", "c", "c", "c"]),
        ("China", "High — loop+guardrails", MINT_MD, ["c", "c", "c", "c"]),
        ("UK", "Medium", PAPER, ["x", "d", "c", "c"]),
        ("US", "Lowest — features only", PAPER, ["x", "d", "c", "d"]),
        ("Germany / EU", "Medium — as warning", PAPER, ["x", "d", "x", "x"]),
    ]
    rh = (5.45 - (y0 + hh)) / len(rows)
    yy = y0 + hh
    for ri, (region, strength, fill, glyphs) in enumerate(rows):
        lc = draw_rect(slide, x0, yy, label_w, rh, fill=fill, line=LINE,
                       line_w=Pt(0.75))
        scol = GREEN_DK if ri < 3 else SLATE
        draw_label(lc, [[(region, {"size": 10.5, "bold": True, "color": INK})],
                        [(strength, {"size": 9, "bold": True, "color": scol})]],
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margins=0.12)
        for c, g in enumerate(glyphs):
            cellx = x0 + label_w + c * cw
            draw_rect(slide, cellx, yy, cw, rh,
                      fill=WHITE if ri % 2 == 0 else CARD_BG, line=LINE,
                      line_w=Pt(0.5))
            kind = {"c": "check", "x": "cross", "d": "dash"}[g]
            tick_glyph(slide, cellx + cw / 2, yy + rh / 2, kind=kind,
                       d=min(0.30, rh - 0.08))
        yy += rh
    dia_caption(slide, caption or ("Borrow mechanisms, reject funding architectures · "
                "all figures company/agency-reported — directional."), y=5.52)


def dia_au_token_loop(slide, caption=None):
    """Circular 4-node token loop (LIME closing arrow) + parallel ASL lane."""
    ccx, ccy = 6.0, 3.45
    nw, nh = 2.4, 0.78
    # node centres on a circle
    top = (ccx, 2.5)
    right = (8.4, 3.45)
    bottom = (ccx, 4.4)
    left = (3.6, 3.45)
    chip(slide, top[0] - nw / 2, top[1] - nh / 2, nw, nh, "Token issued",
         sub="SMS / email · no PIN", fill=GREEN, txt=WHITE, size=11, sub_size=9)
    chip(slide, right[0] - nw / 2, right[1] - nh / 2, nw, nh, "Scan at any pharmacy",
         sub="any phone · forwardable", fill=MINT, txt=GREEN_DK, size=10.5,
         sub_size=9)
    chip(slide, bottom[0] - nw / 2, bottom[1] - nh / 2, nw, nh, "Dispense",
         fill=MINT, txt=GREEN_DK, size=11)
    chip(slide, left[0] - nw / 2, left[1] - nh / 2, nw, nh, "New token per repeat",
         sub="anti-double-dispense built in", fill=MINT_MD, txt=GREEN_DK,
         size=10.5, sub_size=9)
    # clockwise arrows: top→right, right→bottom, bottom→left
    connector(slide, top[0] + nw / 2 - 0.3, top[1] + 0.15,
              right[0] - 0.1, right[1] - nh / 2, color=GREEN, width=1.5)
    connector(slide, right[0] - 0.1, right[1] + nh / 2,
              bottom[0] + nw / 2 - 0.3, bottom[1] - 0.15, color=GREEN, width=1.5)
    connector(slide, bottom[0] - nw / 2 + 0.3, bottom[1] - 0.15,
              left[0] + 0.1, left[1] + nh / 2, color=GREEN, width=1.5)
    # closing arrow left→top — the LIME focal highlight
    connector(slide, left[0] + 0.1, left[1] - nh / 2,
              top[0] - nw / 2 + 0.3, top[1] + 0.15, color=LIME, width=2.5)
    dia_text(slide, ccx - 1.0, ccy - 0.16, 2.0, 0.34,
             [[("e-script token", {"size": 10, "bold": True, "color": SLATE})]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # ASL lane
    dia_text(slide, 0.6, 5.05, 1.7, 0.6,
             [[("Parallel option:", {"size": 9, "bold": True, "color": GREEN_DK})]],
             anchor=MSO_ANCHOR.MIDDLE)
    asl = ["Register once (Active Script List)", "consent per pharmacy (revocable)",
           "all scripts in one managed list"]
    ax = 2.3
    aw = 3.25
    for i, lab in enumerate(asl):
        chip(slide, ax + i * (aw + 0.15), 5.02, aw, 0.55, lab, fill=CARD_BG,
             txt=SLATE, line=LINE, line_w=Pt(0.75), size=9, bold=False)
        if i:
            connector(slide, ax + i * (aw + 0.15) - 0.14, 5.3,
                      ax + i * (aw + 0.15) - 0.01, 5.3, color=SLATE, width=1.0)
    dia_caption(slide, caption or ("Offer token AND list — let the patient choose · "
                "SA needs a printed-QR / USSD / WhatsApp fallback."), y=5.68)


def dia_teleconsult_loop(slide, caption=None, active=None, active_label=None):
    """Left→right 4-stage pipeline with AMBER guardrail chips below.

    `active` (optional) names the stage to emphasise for a given market:
      'ai-triage'  — China-AI: AI-triage guardrail emphasis (1st stage)
      'clinician'  — human-clinician emphasis (2nd stage)
      'pharmacist' — China: pharmacist-review emphasis (3rd stage)
      'order'      — India: order/deliver emphasis (last stage)
    The active stage gets a LIME top-rule + a small caret tag above it.
    `active_label` overrides the caret tag text."""
    stages = [("AI triage / intake", "structured questions", MINT, GREEN_DK),
              ("Human clinician", "diagnoses & e-signs the Rx", GREEN, WHITE),
              ("Pharmacist review", "clinical check", GREEN, WHITE),
              ("Dispense & deliver", "to collect or door", MINT, GREEN_DK)]
    active_idx = {"ai-triage": 0, "clinician": 1, "pharmacist": 2,
                  "order": 3}.get(str(active or "").strip().lower())
    if active_label is None:
        active_label = {0: "AI triage only", 1: "clinician signs",
                        2: "non-delegable", 3: "store-as-node"}.get(active_idx,
                                                                    "active here")
    x = 0.7
    nw, nh, gap = 2.6, 1.2, 0.42
    y = 2.78
    xs = []
    for i, (lab, sub, fill, txt) in enumerate(stages):
        xs.append(x)
        chip(slide, x, y, nw, nh, lab, sub=sub, fill=fill, txt=txt, size=11,
             sub_size=9)
        if i < len(stages) - 1:
            chevron(slide, x + nw + 0.04, y + nh / 2 - 0.18, gap - 0.08, 0.36,
                    fill=GREEN)
        x += nw + gap
    if active_idx is not None:
        ax = xs[active_idx]
        # LIME emphasis rule across the active node + a caret tag above it
        draw_rect(slide, ax, y - 0.02, nw, 0.1, fill=LIME)
        tag = dia_rect(slide, ax + nw / 2 - 0.95, y - 0.42, 1.9, 0.32, fill=LIME,
                       shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        draw_label(tag, [[(str(active_label or "active here"),
                           {"size": 9, "bold": True, "color": GREEN_DK})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.04)
    # GUARDRAILS band
    dia_text(slide, 0.7, 4.0, 12.0, 0.3,
             [[("GUARDRAILS", {"size": 9, "bold": True, "color": GREEN_DK})]])
    guards = [
        (0.7, 3.4, "⛔ No AI-generated prescriptions"),
        (4.42, 3.4, "🛡 No first-visit-online — follow-up / chronic only"),
        (8.14, 3.4, "🛡 Pharmacist review is non-delegable"),
    ]
    for gx, gw, lab in guards:
        box = chip(slide, gx, 4.3, gw, 0.7, lab, fill=PAPER, txt=AMBER,
                   line=AMBER, line_w=Pt(1.25), size=9)
        _line_xml(box, AMBER, 1.25, dashed=True)
    dia_caption(slide, caption or ("China bans AI-generated Rx and enforces "
                "'prescription first, drug second' — a defensible AI/safety stance "
                "for SA."), y=5.4)


def dia_two_economies_one_spine(slide, caption=None):
    """Two equal parallel tracks (insured/cash) converge into one GREEN spine."""
    # Top track — INSURED
    dia_rect(slide, 0.7, 2.4, 8.4, 1.0, fill=MINT, line=GREEN, line_w=Pt(1.25),
             shape=MSO_SHAPE.RECTANGLE)
    chip(slide, 0.82, 2.5, 2.0, 0.4, "INSURED · ~16%", fill=None, txt=GREEN_DK,
         size=11, shape=MSO_SHAPE.RECTANGLE, align=PP_ALIGN.LEFT)
    ins = ["Scheme adjudication", "DSP / formulary", "Co-pay surfaced",
           "Courier / collect"]
    sw = 1.72
    for i, lab in enumerate(ins):
        chip(slide, 0.95 + i * (sw + 0.18), 2.9, sw, 0.46, lab, fill=WHITE,
             txt=GREEN_DK, size=9)
        if i:
            connector(slide, 0.95 + i * (sw + 0.18) - 0.17, 3.13,
                      0.95 + i * (sw + 0.18) - 0.02, 3.13, color=GREEN, width=1.0)
    # Bottom track — CASH (equal weight)
    dia_rect(slide, 0.7, 3.7, 8.4, 1.0, fill=MINT_MD, line=GREEN_DK, line_w=Pt(1.25),
             shape=MSO_SHAPE.RECTANGLE)
    chip(slide, 0.82, 3.8, 2.4, 0.4, "CASH / UNINSURED · ~84%", fill=None,
         txt=GREEN_DK, size=10.5, shape=MSO_SHAPE.RECTANGLE, align=PP_ALIGN.LEFT)
    cash = ["SEP + dispensing fee", "CCMDD / retail pickup", "Price-sensitive choice",
            "Collect / delivery"]
    for i, lab in enumerate(cash):
        chip(slide, 0.95 + i * (sw + 0.18), 4.2, sw, 0.46, lab, fill=WHITE,
             txt=GREEN_DK, size=9)
        if i:
            connector(slide, 0.95 + i * (sw + 0.18) - 0.17, 4.43,
                      0.95 + i * (sw + 0.18) - 0.02, 4.43, color=GREEN_DK, width=1.0)
    # convergence arrows into spine
    connector(slide, 9.1, 2.9, 9.6, 3.3, color=GREEN, width=1.75)
    connector(slide, 9.1, 4.2, 9.6, 3.8, color=GREEN, width=1.75)
    # spine bar
    spine = dia_rect(slide, 9.65, 2.4, 2.95, 2.3, fill=GREEN)
    draw_rect(slide, 9.65, 2.4, 2.95, 0.12, fill=LIME)
    draw_label(spine, [[("ONE SERVICE SPINE", {"size": 12.5, "bold": True,
                                               "color": WHITE})],
                       [("Unified profile", {"size": 9.5, "color": MINT})],
                       [("Omnichannel — same steps, both journeys",
                         {"size": 9.5, "color": MINT})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.12)
    dia_caption(slide, caption or ("~16% insured figure is volatile (state/basis/year) "
                "· data cost ~R20–R79/GB shapes the cash journey."), y=4.92)


def dia_fragmentation_vs_orchestration(slide, caption=None):
    """Signature hero: LEFT tangled fragmentation vs RIGHT clean orchestration."""
    # ---- LEFT HALF (muted, tangled) ----
    dia_text(slide, 0.6, 2.4, 5.7, 0.3,
             [[("TODAY — patient coordinates everything",
                {"size": 10, "bold": True, "color": SLATE})]])
    px, py = 1.2, 3.9
    person_glyph(slide, px, py, scale=0.5, color=SLATE, label="Patient")
    silos = [("Prescriber", 2.6, 2.75), ("Scheme", 4.4, 3.0),
             ("Pharmacy", 4.7, 3.95), ("Stock", 3.9, 4.75),
             ("Delivery", 2.3, 4.95)]
    for lab, sx, sy in silos:
        chip(slide, sx, sy, 1.5, 0.55, lab, fill=PAPER, txt=SLATE, line=LINE,
             line_w=Pt(0.75), size=9.5, bold=False)
        # two-way patient↔silo connector
        connector(slide, px + 0.1, py, sx + 0.1, sy + 0.27, color=SLATE,
                  width=0.9, two_way=True)
    # a couple of broken silo-to-silo dashed connectors (gaps)
    c = connector(slide, 4.1, 3.0, 4.7, 3.95, color=LINE, width=0.9, arrow=False)
    _line_xml(c, LINE, 0.9, dashed=True)
    c = connector(slide, 5.4, 4.2, 4.6, 4.75, color=LINE, width=0.9, arrow=False)
    _line_xml(c, LINE, 0.9, dashed=True)
    # friction micro-labels
    fr = [("re-enter details", 1.62, 3.18), ("chase status", 3.32, 3.5),
          ("call to check", 3.5, 4.5), ("surprise co-pay", 1.35, 4.58)]
    for lab, fx, fy in fr:
        dia_text(slide, fx, fy, 1.85, 0.24,
                 [[(lab, {"size": 9, "italic": True, "color": SLATE})]])
    # ---- CENTRE divider ----
    draw_rect(slide, 6.45, 2.45, 0.014, 3.0, fill=LINE)
    chevron(slide, 6.2, 3.75, 0.5, 0.42, fill=GREEN)
    dia_text(slide, 5.7, 4.2, 1.55, 0.5,
             [[("absorb the coordination", {"size": 9, "italic": True,
                                            "color": GREEN_DK})]],
             align=PP_ALIGN.CENTER)
    # ---- RIGHT HALF (accented, clean) ----
    dia_text(slide, 6.6, 2.4, 6.1, 0.3,
             [[("FUTURE — one system orchestrates",
                {"size": 10, "bold": True, "color": GREEN_DK})]])
    orch = dia_rect(slide, 8.3, 2.7, 3.6, 1.2, fill=GREEN)
    draw_rect(slide, 8.3, 2.7, 3.6, 0.1, fill=LIME)
    draw_label(orch, [[("Treatment Orchestration", {"size": 12, "bold": True,
                                                    "color": WHITE})],
                      [("one system absorbs the coordination",
                        {"size": 9, "color": MINT})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.1)
    # five demoted silos in a neat row
    row_silos = ["Prescriber", "Scheme", "Pharmacy", "Stock", "Delivery"]
    rsx, rsw = 6.7, 1.12
    for i, lab in enumerate(row_silos):
        x = rsx + i * (rsw + 0.05)
        chip(slide, x, 4.35, rsw, 0.5, lab, fill=MINT, txt=GREEN_DK, size=9)
        connector(slide, x + rsw / 2, 4.35, x + rsw / 2, 3.95, color=GREEN,
                  width=1.0)
    # patient + phone mock
    person_glyph(slide, 12.15, 3.05, scale=0.42, color=GREEN_DK, label="Patient",
                 label_color=GREEN_DK)
    connector(slide, 11.9, 3.3, 11.55, 3.3, color=GREEN, width=1.5)
    mock = dia_rect(slide, 9.7, 4.95, 3.0, 0.85, fill=WHITE, line=GREEN,
                    line_w=Pt(1.25))
    draw_label(mock, [[("Status: Ready for collection", {"size": 9, "bold": True,
                                                         "color": INK})],
                      [("Next: Tap to confirm delivery", {"size": 9,
                                                          "color": GREEN_DK})]],
               align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margins=0.12)
    dia_text(slide, 6.7, 5.82, 6.0, 0.3,
             [[("Patient sees ONE state + ONE next step.",
                {"size": 9, "bold": True, "color": GREEN_DK})]])
    dia_text(slide, 0.6, 5.82, 6.0, 0.3,
             [[(caption or "S6 — from dispensing a script to orchestrating the whole "
                "treatment journey.", {"size": 9, "italic": True,
                                       "color": SLATE})]])


def dia_variables_upfront(slide, caption=None):
    """Radial hub-and-spoke: central demand node + five satellite variables."""
    hcx, hcy = 6.67, 3.55
    hub = chip(slide, hcx - 1.3, hcy - 0.65, 2.6, 1.3,
               "Give me all the variables upfront", fill=GREEN, txt=WHITE, size=11)
    draw_rect(slide, hcx - 1.3, hcy - 0.65, 2.6, 0.1, fill=LIME)
    sats = [
        ("Cost", "SEP + dispensing fee, before the till", hcx, 2.4),
        ("Funding / cover", "scheme benefit, co-pay", hcx + 3.0, 2.9),
        ("Stock availability", "is it actually here?", hcx + 3.0, 4.5),
        ("Timing / ETA", "when can I have it?", hcx - 3.0, 4.5),
        ("Alternatives", "generic / substitute options", hcx - 3.0, 2.9),
    ]
    sw, sh = 2.5, 0.82
    for lab, sub, scx, scy in sats:
        connector(slide, hcx, hcy, scx, scy + sh / 2, color=GREEN, width=1.25,
                  arrow=False)
        chip(slide, scx - sw / 2, scy, sw, sh, lab, sub=sub, fill=MINT,
             txt=GREEN_DK, size=10, sub_size=9)
    dia_caption(slide, caption or ("Money is the sharpest case, not the whole of it — "
                "never spring a surprise at the till."), y=5.65)


def dia_strategic_pillars_grid(slide, caption=None):
    """Reuse the engine's native pillar-grid visual (WS3 redundancy guidance):
    the S1–S6 cards with S6 highlighted. Draws the SAME card grid as
    layout_pillar_grid rather than a divergent second grid."""
    cards = [
        ("S1", "Serve both economies on one spine",
         "Insured (~16%) and cash (~84%) share one spine; the cash majority is "
         "never a degraded afterthought."),
        ("S2", "Pharmacy as the health front door",
         "Scripting is the entry to an ongoing care relationship; the store+clinic "
         "estate is a strategic asset."),
        ("S3", "Compete on service, not price",
         "SEP + capped fee fix the price; the durable differentiator is the "
         "experience around the drug."),
        ("S4", "Chronic-medicine flywheel = retention engine",
         "Proactive managed chronic repeats are the recurring relationship that "
         "compounds loyalty."),
        ("S5", "Pharmacist at top-of-licence; automate the toil",
         "Automation/central-fill/telepharmacy absorb mechanics so pharmacists do "
         "the clinical work."),
        ("S6", "From fulfilment to treatment orchestration",
         "Absorb the coordination across prescriber, pharmacy, scheme, stock, "
         "delivery into one system."),
    ]
    _pillar_grid(slide, cards, top=Inches(2.4), cols=3, accent=GREEN,
                 card_h=Inches(1.52), small=False)
    # Highlight S6 (last card) with a GREEN border + LIME "NEW" chip.
    lx, gx = Inches(0.6), Inches(0.22)
    avail_w = SW - lx * 2
    cw = (avail_w - gx * 2) / 3
    card_h = Inches(1.52)
    gy = Inches(0.22)
    x = lx + 2 * (cw + gx)
    y = Inches(2.4) + 1 * (card_h + gy)
    border = add_rect(slide, x, y, cw, card_h, fill=None, line=GREEN, line_w=Pt(2),
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    newchip = add_rect(slide, x + cw - Inches(0.75), y + Inches(0.16), Inches(0.6),
                       Inches(0.3), fill=LIME, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(newchip, [[("NEW", {"size": 9, "bold": True, "color": GREEN_DK})]],
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)


def dia_pattern_pillar_map(slide, caption=None):
    """Two-column bipartite map: borrowed patterns (left) → UX pillars (right)."""
    patterns = [
        "QR e-script token, scan anywhere",
        "New token per repeat (anti-double-dispense)",
        "Nominated pharmacy / Active Script List",
        "Real-time order tracking (NHS-App style)",
        "Auto-refill / med-sync default",
        "Dose-by-time packaging + app (PillPack)",
        "Price transparency at point of commit",
    ]
    pillars = [
        "U1 — Channel ladder (WhatsApp-first)",
        "U2 — Never leave them guessing",
        "U3 — All the variables upfront",
        "U4 — Design for graceful failure",
        "U5 — Dignity, privacy & language",
        "U6 — Reduce the journey, don't just digitise",
        "U7 — Adherence-as-design",
    ]
    lx, rx, w, h = 0.7, 8.0, 4.6, 0.5
    y0, gap = 2.45, 0.44
    lcent, rcent = {}, {}
    for i, lab in enumerate(patterns):
        y = y0 + i * (h + gap)
        chip(slide, lx, y, w, h, lab, fill=MINT, txt=GREEN_DK, size=9.5,
             align=PP_ALIGN.LEFT)
        lcent[i + 1] = (lx + w, y + h / 2)
    for i, lab in enumerate(pillars):
        y = y0 + i * (h + gap)
        is_u5 = lab.startswith("U5")
        chip(slide, rx, y, w, h, lab, fill=MINT_MD, txt=GREEN_DK, size=9.5,
             align=PP_ALIGN.LEFT)
        rcent[i + 1] = (rx, y + h / 2)
        if is_u5:
            tag = draw_rect(slide, rx + w - 1.18, y + 0.07, 1.08, 0.36, fill=LIME,
                            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            draw_label(tag, [[("SA-specific", {"size": 9, "bold": True,
                                               "color": GREEN_DK})]],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
    # connectors: (pattern# -> pillar#); 4->2 is the exemplar (GREEN_DK, thicker)
    links = [(1, 6), (1, 1), (2, 4), (3, 6), (4, 2, True), (5, 7), (6, 7), (7, 3)]
    for link in links:
        p, u = link[0], link[1]
        emph = len(link) > 2 and link[2]
        x1, y1 = lcent[p]
        x2, y2 = rcent[u]
        connector(slide, x1, y1, x2, y2,
                  color=GREEN_DK if emph else GREEN,
                  width=2.0 if emph else 1.0, arrow=False)
    dia_caption(slide, caption or ("Borrowed mechanisms map onto the UX constitution; "
                "U5 (dignity/POPIA) is the SA-specific layer no benchmark supplied."),
                y=5.7)


# Native-diagram dispatch (dia:<id> -> drawing fn(slide, caption)).
DIAGRAMS = {
    "front-door": dia_front_door,
    "fulfilment-node": dia_fulfilment_node,
    "time-reallocation": dia_time_reallocation,
    "efficacy-effectiveness": dia_efficacy_effectiveness,
    "script-artefact": dia_script_artefact,
    "transfer-fit": dia_transfer_fit,
    "au-token-loop": dia_au_token_loop,
    "teleconsult-loop": dia_teleconsult_loop,
    "two-economies-one-spine": dia_two_economies_one_spine,
    "fragmentation-vs-orchestration": dia_fragmentation_vs_orchestration,
    "variables-upfront": dia_variables_upfront,
    "strategic-pillars-grid": dia_strategic_pillars_grid,
    "pattern-pillar-map": dia_pattern_pillar_map,
}


def is_dia(visual):
    return bool(visual) and str(visual).strip().startswith("dia:")


def render_diagram(slide, visual, caption=None):
    """Render a native dia:<id> into the FULL diagram canvas. Returns True if a
    real diagram was drawn, False if unknown (caller may fall back)."""
    key = str(visual).strip()[4:]
    fn = DIAGRAMS.get(key)
    if fn is None:
        sys.stderr.write(f"[warn] unknown diagram id: {visual}\n")
        return False
    # A long content caption (e.g. the S6 hero's prose) won't fit the one-line
    # diagram caption slot — let the diagram's own concise caption win.
    if caption and len(str(caption)) > 140:
        caption = None
    try:
        fn(slide, caption=caption)
    except Exception as exc:
        sys.stderr.write(f"[warn] diagram '{visual}' failed: {exc}\n")
        return False
    return True


def _render_visual(slide, visual, x, y, w, h, caption=None):
    """Render a `visual:` id at the given box.
    - img:<id>  -> embed PNG from deck-assets (fit within box, centred).
    - dia:<id>  -> native diagram, drawn full-canvas (box args ignored).
    Returns True if anything was drawn.
    """
    if not visual:
        return False
    visual = str(visual).strip()
    if visual.startswith("img:"):
        key = visual[4:]
        path = IMG_MAP.get(key)
        if path and os.path.exists(path):
            # frame + fit-within (preserve aspect by width, clamp height)
            try:
                from PIL import Image
                iw_px, ih_px = Image.open(path).size
                aspect = iw_px / ih_px if ih_px else 1.524
            except Exception:
                aspect = 1.524
            draw_w = w
            draw_h = Emu(int(int(w) / aspect))
            if int(draw_h) > int(h):
                draw_h = h
                draw_w = Emu(int(int(h) * aspect))
            ox = x + Emu(int((int(w) - int(draw_w)) / 2))
            oy = y + Emu(int((int(h) - int(draw_h)) / 2))
            add_rect(slide, ox - Inches(0.05), oy - Inches(0.05),
                     draw_w + Inches(0.1), draw_h + Inches(0.1),
                     fill=WHITE, line=LINE, line_w=Pt(1), shadow=True)
            slide.shapes.add_picture(path, ox, oy, width=draw_w, height=draw_h)
            if caption:
                add_text(slide, x, oy + draw_h + Inches(0.06), w, Inches(0.3),
                         [[(caption, {"size": 8.5, "italic": True, "color": SLATE})]],
                         align=PP_ALIGN.CENTER)
            return True
        # missing image -> placeholder
        return _visual_placeholder(slide, x, y, w, h,
                                   f"image not found: {visual}", caption)
    if visual.startswith("dia:"):
        # Native diagrams own the full canvas; the per-layout box is ignored.
        if render_diagram(slide, visual, caption=caption):
            return True
        return _visual_placeholder(slide, x, y, w, h,
                                   f"unknown diagram: {visual}", caption)
    # unknown visual id
    return _visual_placeholder(slide, x, y, w, h, f"unknown visual: {visual}", caption)


def _visual_placeholder(slide, x, y, w, h, label, caption=None):
    box = add_rect(slide, x, y, w, h, fill=CARD_BG, line=GREEN,
                   line_w=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # dashed feel via lime corner tab + centred label
    add_rect(slide, x, y, Inches(0.12), h, fill=LIME)
    runs = [[(line, {"size": 12, "bold": True, "color": GREEN_DK})]
            for line in str(label).split("\n")]
    shape_text(box, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.2)
    if caption:
        add_text(slide, x, y + h + Inches(0.06), w, Inches(0.4),
                 [[(caption, {"size": 8.5, "italic": True, "color": SLATE})]],
                 align=PP_ALIGN.CENTER, line_spacing=0.95)
    return True


# ============================================================================
# LAYOUT-TEMPLATE RENDERERS (SPEC §2). Each consumes a `block` dict (SPEC §3).
# ============================================================================

def _g(block, key, default=None):
    """Get a field from a content block, tolerant of missing keys."""
    v = block.get(key, default)
    return v if v is not None else default


def _visual_of(block):
    """Return the slide's visual id, or None if absent/explicitly omitted.
    Content authors may write `visual: omit` (or none/-) to mean 'no visual'."""
    v = _g(block, "visual")
    if not v:
        return None
    s = str(v).strip()
    if s.lower() in ("omit", "none", "-", "n/a", "na", ""):
        return None
    return s


def _render_diagram_block(slide, visual, block, caption=None):
    """Render a native dia:<id> into the FULL diagram canvas (s == 1.0, so labels
    keep their authored sizes — the legible reference, same as slides 18/30).
    Honours an optional per-slide `diagram_active` field for diagrams that take an
    active-stage highlight (currently dia:teleconsult-loop, N2)."""
    key = str(visual).strip()[4:]
    active = _g(block, "diagram_active")
    if key == "teleconsult-loop" and active:
        try:
            dia_teleconsult_loop(slide, caption=caption, active=str(active))
            return True
        except Exception as exc:
            sys.stderr.write(f"[warn] diagram '{visual}' failed: {exc}\n")
            return False
    return render_diagram(slide, visual, caption=caption)


def _emit_diagram_slide(prs, block, visual, title, takeaway):
    """Emit a DEDICATED follow-on slide carrying a native diagram full-canvas.

    B1 fix: a `dia:` paired with a text column (pillar / market-profile) used to be
    squeezed into a narrow side region, collapsing labels to ~4pt. Instead the
    parent slide keeps the text and the diagram gets its own full-width canvas
    here, so every label renders at its authored size (>= the 9pt floor), exactly
    like the proven full-canvas diagrams on slides 18 and 30."""
    slide = new_slide(prs)
    # The so-what band is one strip (~2 lines). The diagram below carries the
    # detail in its own caption, so keep the band's takeaway from overflowing by
    # softly trimming a very long takeaway at a sentence/clause boundary.
    tk = str(takeaway or "")
    if len(tk) > 200:
        cut = tk.rfind(" — ", 0, 200)
        if cut < 80:
            cut = tk.rfind(", ", 0, 200)
        if cut < 80:
            cut = tk.rfind(" ", 0, 197)
        tk = (tk[:cut].rstrip(" ,—") + " …") if cut > 0 else (tk[:197] + "…")
    content_chrome(slide, _g(block, "kicker", ""), title, tk,
                   cite=_evidence_line(_g(block, "evidence")))
    _render_diagram_block(slide, visual, block, caption=_g(block, "caption"))
    _flag_strip(slide, _g(block, "flags"))
    return slide


def layout_cover(prs, block):
    slide = new_slide(prs)
    _set_bg(slide, WHITE)
    add_rect(slide, 0, 0, Inches(0.45), SH, fill=GREEN)
    add_rect(slide, Inches(0.45), 0, Inches(0.08), SH, fill=LIME)
    add_rect(slide, Inches(10.55), Inches(0.9), Inches(0.62), Inches(2.9), fill=MINT_MD)
    add_rect(slide, Inches(9.4), Inches(2.04), Inches(2.9), Inches(0.62), fill=MINT_MD)
    add_text(slide, Inches(1.0), Inches(1.5), Inches(9.0), Inches(0.4),
             [[("BIGLY LABS  ×  DIS-CHEM", {"size": 13, "bold": True, "color": GREEN})]])
    title = _g(block, "title", "Future-State Scripting")
    add_text(slide, Inches(0.95), Inches(2.25), Inches(11.4), Inches(2.2),
             [[(title, {"size": 46, "bold": True, "color": INK})]], line_spacing=1.0)
    subtitle = _g(block, "subtitle", "")
    if subtitle:
        add_text(slide, Inches(1.0), Inches(4.55), Inches(11.0), Inches(0.9),
                 [[(subtitle, {"size": 16, "color": SLATE})]], line_spacing=1.15)
    rb = add_rect(slide, Inches(1.0), Inches(5.5), Inches(8.7), Inches(0.5),
                  fill=MINT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(rb, [[("⚠  Concept for review — not an official Dis-Chem artefact.",
                      {"size": 10.5, "bold": True, "color": GREEN_DK})]],
               anchor=MSO_ANCHOR.MIDDLE, margins=0.14)
    date = _g(block, "date", "")
    foot = "Executive strategy deck for Dis-Chem & Bigly leadership"
    if date:
        foot += "  ·  " + date
    add_text(slide, Inches(1.0), Inches(6.7), Inches(10.0), Inches(0.4),
             [[(foot, {"size": 11, "color": SLATE})]])
    PAGE["n"] += 1
    _maybe_notes(slide, block)
    return slide


def layout_exec_summary(prs, block):
    slide = new_slide(prs)
    content_chrome(slide, _g(block, "kicker", "Executive summary"),
                   _g(block, "title", "The strategy in one slide"),
                   _g(block, "so_what", ""), cite=_evidence_line(_g(block, "evidence")))
    top = Inches(2.45)
    body = _g(block, "body", []) or []
    visual = _visual_of(block)
    if visual:
        # body left, visual right
        colw = Inches(7.0)
        _bullet_block(slide, Inches(0.6), top + Inches(0.1), colw - Inches(0.2),
                      body, size=12.5, gap=0.78)
        _render_visual(slide, visual, Inches(7.9), top, Inches(4.85), Inches(4.0),
                       caption=None)
    else:
        # single emphasis card with bullets
        add_rect(slide, Inches(0.6), top, Inches(12.13), Inches(4.1), fill=CARD_BG,
                 line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        _bullet_block(slide, Inches(0.95), top + Inches(0.35), Inches(11.4),
                      body, size=14, gap=0.82)
    _flag_strip(slide, _g(block, "flags"))
    _maybe_notes(slide, block)
    return slide


def layout_section_divider(prs, block):
    return section_divider(prs, str(_g(block, "num", "")),
                            _g(block, "title", ""), _g(block, "subtitle", ""))


def layout_theme(prs, block):
    slide = new_slide(prs)
    content_chrome(slide, _g(block, "kicker", ""), _g(block, "title", "Untitled"),
                   _g(block, "so_what", ""), cite=_evidence_line(_g(block, "evidence")))
    top = Inches(2.45)
    body = _g(block, "body", []) or []
    visual = _visual_of(block)
    if is_dia(visual):
        # Native diagram owns the full canvas; bullets are summarised by the so-what.
        _render_diagram_block(slide, visual, block)
    elif visual:
        colw = Inches(6.7)
        _bullet_block(slide, Inches(0.6), top, colw - Inches(0.2), body,
                      size=12.5, gap=0.74)
        _render_visual(slide, visual, Inches(7.55), top, Inches(5.2), Inches(3.9),
                       caption=None)
    else:
        _bullet_block(slide, Inches(0.6), top, Inches(8.2), body, size=13.5, gap=0.8)
        ev = _g(block, "evidence")
        # evidence pulled into the cite line already; nothing extra needed
    _flag_strip(slide, _g(block, "flags"))
    _maybe_notes(slide, block)
    return slide


def layout_market_profile(prs, block):
    slide = new_slide(prs)
    content_chrome(slide, _g(block, "kicker", "Global benchmark"),
                   _g(block, "title", "Untitled market"),
                   _g(block, "so_what", ""), cite=_evidence_line(_g(block, "evidence")))
    top = Inches(2.45)
    # Left: fact tiles. Right: body bullets (+ optional visual below body).
    facts = _g(block, "facts", []) or []
    lx = Inches(0.6)
    lw = Inches(5.4)
    add_text(slide, lx, top, lw, Inches(0.3),
             [[("KEY FACTS", {"size": 11, "bold": True, "color": GREEN})]])
    yy = top + Inches(0.4)
    for fact in facts[:5]:
        label = fact.get("label", "") if isinstance(fact, dict) else ""
        value = fact.get("value", "") if isinstance(fact, dict) else str(fact)
        card = add_rect(slide, lx, yy, lw, Inches(0.72), fill=WHITE, line=LINE,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        add_rect(slide, lx, yy, Inches(0.1), Inches(0.72), fill=GREEN)
        add_text(slide, lx + Inches(0.28), yy + Inches(0.08), lw - Inches(0.5), Inches(0.3),
                 [[(label, {"size": 10, "bold": True, "color": SLATE})]])
        add_text(slide, lx + Inches(0.28), yy + Inches(0.36), lw - Inches(0.5), Inches(0.32),
                 [[(value, {"size": 13, "bold": True, "color": GREEN_DK})]])
        yy += Inches(0.82)
    # Right: body, OR (for a native diagram) the loop is emitted on its own slide.
    rx = Inches(6.25)
    rw = Inches(6.48)
    body = _g(block, "body", []) or []
    visual = _visual_of(block)
    if is_dia(visual):
        # B1: do NOT cram the loop into a narrow side region (it collapsed labels
        # to ~4pt). Facts stay left, body bullets fill the right column here, and
        # the diagram gets its own full-width slide next.
        add_text(slide, rx, top, rw, Inches(0.3),
                 [[("HOW THE LOOP RUNS", {"size": 11, "bold": True, "color": GREEN})]])
        _bullet_block(slide, rx, top + Inches(0.4), rw, body, size=11.5, gap=0.62)
        _flag_strip(slide, _g(block, "flags"))
        _maybe_notes(slide, block)
        dtitle = f"{_g(block, 'title', 'The loop')} — how the loop runs"
        _emit_diagram_slide(prs, block, visual, dtitle, _g(block, "so_what", ""))
        return slide
    elif visual:
        add_text(slide, rx, top, rw, Inches(0.3),
                 [[("WHAT IT PROVES", {"size": 11, "bold": True, "color": GREEN})]])
        _bullet_block(slide, rx, top + Inches(0.4), rw, body, size=11.5, gap=0.62)
        _render_visual(slide, visual, rx, top + Inches(2.7), rw, Inches(1.9), caption=None)
    else:
        add_text(slide, rx, top, rw, Inches(0.3),
                 [[("WHAT IT PROVES", {"size": 11, "bold": True, "color": GREEN})]])
        _bullet_block(slide, rx, top + Inches(0.4), rw, body, size=12.5, gap=0.7)
    _flag_strip(slide, _g(block, "flags"))
    _maybe_notes(slide, block)
    return slide


def layout_comparison(prs, block):
    slide = new_slide(prs)
    content_chrome(slide, _g(block, "kicker", ""),
                   _g(block, "title", "Comparison"),
                   _g(block, "so_what", ""), cite=_evidence_line(_g(block, "evidence")))
    top = Inches(2.45)
    visual = _visual_of(block)
    # A native diagram on a comparison slide (e.g. dia:pattern-pillar-map) is the
    # intended visual per diagrams.md — render it full-canvas instead of a table.
    if is_dia(visual):
        render_diagram(slide, visual)
        _flag_strip(slide, _g(block, "flags"))
        _maybe_notes(slide, block)
        return slide
    columns = _g(block, "columns", []) or []
    rows = _g(block, "rows", []) or []
    lx = Inches(0.6)
    table_w = SW - lx * 2
    ncols = max(len(columns), 1)
    cw = Emu(int(table_w) // ncols)
    # header row
    hh = Inches(0.5)
    for c, col in enumerate(columns):
        x = lx + Emu(int(cw) * c)
        cell = add_rect(slide, x, top, cw, hh, fill=GREEN, line=WHITE, line_w=Pt(1.5))
        shape_text(cell, [[(str(col), {"size": 11.5, "bold": True, "color": WHITE})]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.08)
    # body rows
    rh = Inches(0.56)
    yy = top + hh
    for ri, row in enumerate(rows):
        cells = row if isinstance(row, list) else [row]
        bg = WHITE if ri % 2 == 0 else CARD_BG
        for c in range(ncols):
            x = lx + Emu(int(cw) * c)
            val = str(cells[c]) if c < len(cells) else ""
            cell = add_rect(slide, x, yy, cw, rh, fill=bg, line=LINE, line_w=Pt(0.75))
            bold = (c == 0)
            col_txt = INK if c == 0 else SLATE
            shape_text(cell, [[(val, {"size": 10.5, "bold": bold, "color": col_txt})]],
                       align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE, margins=0.12)
        yy += rh
    _flag_strip(slide, _g(block, "flags"))
    _maybe_notes(slide, block)
    return slide


def layout_pillar(prs, block):
    slide = new_slide(prs)
    pid = str(_g(block, "id", ""))
    name = _g(block, "name", "Strategic pillar")
    title = f"{pid} — {name}" if pid else name
    content_chrome(slide, _g(block, "kicker", "Strategic pillar"), title,
                   _g(block, "what", ""), cite=_evidence_line(_g(block, "evidence")))
    top = Inches(2.45)
    visual = _visual_of(block)
    dia = is_dia(visual)
    # B1: a native diagram no longer shares the slide as a cramped side box; the
    # text uses the full width and the diagram is emitted on its own slide next.
    colw = Inches(6.6) if (visual and not dia) else Inches(12.13)
    # WHY card
    why = _g(block, "why", "")
    add_text(slide, Inches(0.6), top, colw, Inches(0.3),
             [[("WHY THIS MATTERS", {"size": 11, "bold": True, "color": GREEN})]])
    add_text(slide, Inches(0.6), top + Inches(0.4), colw, Inches(1.6),
             [[(why, {"size": 13, "color": INK})]], line_spacing=1.12)
    # EXAMPLE band
    example = _g(block, "example", "")
    if example:
        ey = top + Inches(2.2)
        band = add_rect(slide, Inches(0.6), ey, colw, Inches(1.5), fill=MINT,
                        line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        shape_text(band, [[("In practice:  ", {"size": 12, "bold": True, "color": GREEN_DK}),
                           (example, {"size": 12, "color": INK})]],
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, margins=0.2)
    if dia:
        _flag_strip(slide, _g(block, "flags"))
        _maybe_notes(slide, block)
        # Clean title for the dedicated diagram slide: prefer the block's real
        # `title` (or `name`), avoiding the slug-prefixed pillar header.
        base = _g(block, "title") or _g(block, "name") or name
        dtitle = f"{base} — the pattern"
        # Takeaway for the diagram slide: prefer the concise `so_what`/`what`
        # (not the longer `why`, which can overflow the one-line so-what band).
        dtake = _g(block, "so_what", "") or _g(block, "what", "")
        _emit_diagram_slide(prs, block, visual, dtitle, dtake)
        return slide
    elif visual:
        _render_visual(slide, visual, Inches(7.45), top, Inches(5.3), Inches(4.0),
                       caption=None)
    _flag_strip(slide, _g(block, "flags"))
    _maybe_notes(slide, block)
    return slide


def layout_pillar_grid(prs, block):
    slide = new_slide(prs)
    content_chrome(slide, _g(block, "kicker", ""),
                   _g(block, "title", "Strategic pillars"),
                   _g(block, "so_what", ""), cite=_evidence_line(_g(block, "evidence")))
    cards = _g(block, "cards", []) or []
    items = []
    for card in cards:
        if isinstance(card, dict):
            items.append((str(card.get("id", "")), card.get("name", ""),
                          card.get("oneLiner", card.get("oneliner", ""))))
        else:
            items.append(("", str(card), ""))
    n = len(items)
    cols = 3 if n <= 6 else 4
    # size cards to fit the count
    rows = (n + cols - 1) // cols if n else 1
    card_h = Inches(1.72) if rows <= 2 else Inches(1.42)
    top = Inches(2.45)
    _pillar_grid(slide, items, top=top, cols=cols, accent=GREEN,
                 card_h=card_h, small=(n > 6))
    # WS3 redundancy guidance: `dia:strategic-pillars-grid` is the SAME visual as
    # this native grid — do NOT double-draw it. We instead apply the diagram's
    # one differentiator here: highlight S6 (GREEN border + LIME "NEW" chip).
    visual = _visual_of(block)
    if is_dia(visual):
        s6 = next((i for i, it in enumerate(items)
                   if str(it[0]).upper() == "S6"), None)
        if s6 is not None:
            gx = Inches(0.22)
            lxx = Inches(0.6)
            cw = (SW - lxx * 2 - gx * (cols - 1)) / cols
            gy = Inches(0.22)
            r, c = divmod(s6, cols)
            x = lxx + c * (cw + gx)
            y = top + r * (card_h + gy)
            add_rect(slide, x, y, cw, card_h, fill=None, line=GREEN, line_w=Pt(2),
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            ch = add_rect(slide, x + cw - Inches(0.78), y + Inches(0.16),
                          Inches(0.6), Inches(0.3), fill=LIME,
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            shape_text(ch, [[("NEW", {"size": 9, "bold": True, "color": GREEN_DK})]],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margins=0.02)
    _flag_strip(slide, _g(block, "flags"))
    _maybe_notes(slide, block)
    return slide


def layout_diagram_hero(prs, block):
    slide = new_slide(prs)
    content_chrome(slide, _g(block, "kicker", ""),
                   _g(block, "title", "Diagram"),
                   _g(block, "so_what", ""), cite=_evidence_line(_g(block, "evidence")))
    top = Inches(2.4)
    visual = _visual_of(block)
    caption = _g(block, "caption", "")
    _render_visual(slide, visual, Inches(1.6), top, Inches(10.13), Inches(3.9),
                   caption=caption)
    _flag_strip(slide, _g(block, "flags"))
    _maybe_notes(slide, block)
    return slide


def layout_sources(prs, block):
    slide = new_slide(prs)
    content_chrome(slide, _g(block, "kicker", "Appendix"),
                   _g(block, "title", "Sources & evidence"),
                   _g(block, "so_what", ""))
    items = _g(block, "items", []) or _g(block, "body", []) or []
    top_in = 2.4
    bottom_in = 7.04           # stay above the footer hairline (7.18)
    # two-column list if long
    half = (len(items) + 1) // 2
    cols = [items[:half], items[half:]] if len(items) > 8 else [items]
    colw = Inches(6.0) if len(cols) == 2 else Inches(12.0)
    per_col = max((len(c) for c in cols), default=1)
    pitch = min(0.36, (bottom_in - top_in) / max(per_col, 1))
    size = 10.5 if per_col <= 11 else (9.5 if per_col <= 13 else 8.7)
    for ci, col_items in enumerate(cols):
        x = Inches(0.6) + ci * (colw + Inches(0.25))
        yy = top_in
        for it in col_items:
            add_text(slide, x, Inches(yy), colw, Inches(pitch),
                     [[("• ", {"size": size - 0.5, "bold": True, "color": GREEN}),
                       (str(it), {"size": size, "color": INK})]], line_spacing=0.98)
            yy += pitch
    _maybe_notes(slide, block)
    return slide


def layout_closing(prs, block):
    slide = new_slide(prs)
    content_chrome(slide, _g(block, "kicker", "What's next"),
                   _g(block, "title", "Where we go from here"),
                   _g(block, "so_what", ""), cite=_evidence_line(_g(block, "evidence")))
    body = _g(block, "body", []) or []
    flags = _g(block, "flags")
    top_in = 2.5
    # The card must fit between the so-what band and the flag strip / footer.
    bottom_in = 6.62 if flags else 6.92
    card_h = bottom_in - top_in
    add_rect(slide, Inches(0.6), Inches(top_in), Inches(12.13), Inches(card_h),
             fill=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    n = max(len(body), 1)
    pad = 0.34
    avail = card_h - 2 * pad
    gap = min(0.85, avail / n)               # adaptive row pitch
    # font scales down a touch for dense lists so each line fits its row
    size = 15 if n <= 5 else (13 if n <= 6 else (11.5 if n <= 7 else 10.5))
    row_h = max(0.34, gap)
    yy = top_in + pad
    for it in body:
        add_text(slide, Inches(1.0), Inches(yy), Inches(11.3), Inches(row_h),
                 [[("▸ ", {"size": size, "bold": True, "color": LIME}),
                   (str(it), {"size": size, "color": WHITE})]], line_spacing=1.0)
        yy += gap
    _flag_strip(slide, flags)
    _maybe_notes(slide, block)
    return slide


def _maybe_notes(slide, block):
    nt = _g(block, "notes")
    if nt:
        try:
            notes(slide, str(nt))
        except Exception:
            pass


# Layout dispatch table (SPEC §2 catalogue).
LAYOUTS = {
    "cover": layout_cover,
    "exec-summary": layout_exec_summary,
    "section-divider": layout_section_divider,
    "theme": layout_theme,
    "market-profile": layout_market_profile,
    "comparison": layout_comparison,
    "pillar": layout_pillar,
    "pillar-grid": layout_pillar_grid,
    "diagram-hero": layout_diagram_hero,
    "sources": layout_sources,
    "closing": layout_closing,
}


def _warning_slide(prs, message, detail=""):
    """A clearly-labelled slide for unknown layouts / fatal block errors."""
    slide = new_slide(prs)
    content_chrome(slide, "Build warning", message, detail or "")
    add_rect(slide, Inches(0.6), Inches(2.6), Inches(12.13), Inches(2.5),
             fill=RGBColor(0xFD, 0xF6, 0xE8), line=RGBColor(0xE6, 0xCB, 0x8A),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, Inches(1.0), Inches(2.9), Inches(11.3), Inches(2.0),
             [[("⚠  ", {"size": 16, "bold": True, "color": AMBER}),
               (message, {"size": 14, "bold": True, "color": INK})],
              [(detail, {"size": 11, "color": SLATE})]], line_spacing=1.2)
    return slide


# ============================================================================
# TOLERANT CONTENT LOADER (SPEC §3)
# ============================================================================
# Format per part file: an ordered list of slide blocks, each opened by a line
#   ## SLIDE: <layout>
# followed by `key: value` lines. Multi-line list fields are written as:
#   body:
#     - item one
#     - item two
# The `facts:` field is a list of inline {label: .., value: ..} dicts:
#   facts:
#     - {label: "Apollo 24|7 orders", value: "~75k/day"}
# `evidence:`, `flags:`, `columns:` may be inline JSON-ish arrays  [a, b]  OR a
# multi-line `- item` list. `rows:` (comparison) is a multi-line list whose
# items may be inline arrays. Unknown layouts/fields never crash the build.

LIST_FIELDS = {"body", "evidence", "flags", "items", "columns", "cards", "rows", "facts"}


def _strip_inline_comment(s):
    # remove a trailing  "# ..."  comment that is not inside quotes/braces
    out = []
    in_s = in_d = False
    depth = 0
    for ch in s:
        if ch == "'" and not in_d:
            in_s = not in_s
        elif ch == '"' and not in_s:
            in_d = not in_d
        elif ch in "{[" and not in_s and not in_d:
            depth += 1
        elif ch in "}]" and not in_s and not in_d:
            depth = max(0, depth - 1)
        elif ch == "#" and not in_s and not in_d and depth == 0:
            break
        out.append(ch)
    return "".join(out).rstrip()


def _unquote(s):
    s = s.strip()
    if len(s) >= 2 and ((s[0] == s[-1] == '"') or (s[0] == s[-1] == "'")):
        return s[1:-1]
    return s


def _parse_inline_array(s):
    """Parse a bracketed inline array like [wiki/x, src-y] -> list of strings."""
    s = s.strip()
    if s.startswith("[") and s.endswith("]"):
        s = s[1:-1]
    if not s.strip():
        return []
    # split on commas not inside braces/brackets
    items, depth, cur = [], 0, ""
    for ch in s:
        if ch in "{[":
            depth += 1
        elif ch in "}]":
            depth -= 1
        if ch == "," and depth == 0:
            items.append(cur)
            cur = ""
        else:
            cur += ch
    if cur.strip():
        items.append(cur)
    return [_unquote(i) for i in items if i.strip()]


def _parse_fact(s):
    """Parse  {label: "x", value: "y"}  -> {'label': 'x', 'value': 'y'}."""
    s = s.strip()
    if s.startswith("{") and s.endswith("}"):
        s = s[1:-1]
    out = {}
    # split top-level commas
    parts, depth, cur = [], 0, ""
    in_s = in_d = False
    for ch in s:
        if ch == "'" and not in_d:
            in_s = not in_s
        elif ch == '"' and not in_s:
            in_d = not in_d
        elif ch in "{[" and not in_s and not in_d:
            depth += 1
        elif ch in "}]" and not in_s and not in_d:
            depth -= 1
        if ch == "," and depth == 0 and not in_s and not in_d:
            parts.append(cur)
            cur = ""
        else:
            cur += ch
    if cur.strip():
        parts.append(cur)
    for p in parts:
        if ":" in p:
            k, _, v = p.partition(":")
            out[k.strip()] = _unquote(v)
    return out


def _parse_list_item(field, raw):
    """Convert a single `- ...` list item for the given field."""
    raw = raw.strip()
    if field == "facts":
        return _parse_fact(raw)
    if field == "rows" and raw.startswith("["):
        return _parse_inline_array(raw)
    if raw.startswith("{") and field != "rows":
        return _parse_fact(raw)
    return _unquote(raw)


def parse_part_file(path):
    """Parse one part*.md into an ordered list of block dicts. Never raises."""
    blocks = []
    try:
        with open(path, "r", encoding="utf-8") as fh:
            lines = fh.read().splitlines()
    except Exception as exc:
        sys.stderr.write(f"[warn] could not read {path}: {exc}\n")
        return blocks

    cur = None
    cur_list_field = None  # name of a multi-line list field we are accumulating
    header_re = re.compile(r"^##\s*SLIDE:\s*(.+?)\s*$", re.IGNORECASE)
    keyval_re = re.compile(r"^([A-Za-z_][\w-]*)\s*:\s*(.*)$")

    def _close():
        if cur is not None:
            blocks.append(cur)

    for raw_line in lines:
        line = raw_line.rstrip("\n")
        m = header_re.match(line.strip())
        if m:
            _close()
            cur = {"_layout": m.group(1).strip().lower(), "_src": os.path.basename(path)}
            cur_list_field = None
            continue
        if cur is None:
            continue  # ignore preamble before first SLIDE

        stripped = line.strip()
        # multi-line list item?
        if cur_list_field and (stripped.startswith("- ") or stripped == "-"):
            item_raw = stripped[1:].strip()
            try:
                cur.setdefault(cur_list_field, []).append(
                    _parse_list_item(cur_list_field, item_raw))
            except Exception as exc:
                sys.stderr.write(f"[warn] {path}: bad list item in "
                                 f"'{cur_list_field}': {item_raw!r} ({exc})\n")
            continue

        if stripped == "":
            # blank line does not end a list block (allows spacing); skip
            continue

        km = keyval_re.match(line)
        if km:
            key = km.group(1).strip().lower()
            val = _strip_inline_comment(km.group(2)).strip()
            cur_list_field = None
            if val == "":
                # could be the opener for a multi-line list field
                if key in LIST_FIELDS:
                    cur_list_field = key
                    cur.setdefault(key, [])
                else:
                    cur[key] = ""
                continue
            # inline value
            if key in LIST_FIELDS and val.startswith("["):
                cur[key] = _parse_inline_array(val)
            elif key in LIST_FIELDS:
                # single inline scalar for a list field -> wrap
                cur[key] = [_unquote(val)]
            else:
                cur[key] = _unquote(val)
            continue
        # a stray non-key line: append to last scalar-ish (ignore safely)
        # (kept silent to avoid noise on partially-written files)

    _close()
    return blocks


def _part_sort_key(path):
    """Order files by the integer in 'part<N>' then by filename."""
    base = os.path.basename(path)
    m = re.search(r"part(\d+)", base, re.IGNORECASE)
    n = int(m.group(1)) if m else 9999
    return (n, base)


def load_content():
    """Load all part*.md files in order -> flat list of block dicts."""
    if not os.path.isdir(CONTENT_DIR):
        sys.stderr.write(f"[warn] content dir not found: {CONTENT_DIR}\n")
        return []
    files = sorted(glob.glob(os.path.join(CONTENT_DIR, "part*.md")),
                   key=_part_sort_key)
    all_blocks = []
    for f in files:
        bs = parse_part_file(f)
        sys.stderr.write(f"[info] {os.path.basename(f)}: {len(bs)} slide block(s)\n")
        all_blocks.extend(bs)
    return all_blocks


# ============================================================================
# ASSEMBLY
# ============================================================================

def _render_block(prs, block):
    layout = block.get("_layout", "")
    fn = LAYOUTS.get(layout)
    if fn is None:
        sys.stderr.write(f"[warn] unknown layout '{layout}' "
                         f"(from {block.get('_src','?')}) — rendering warning slide\n")
        _warning_slide(prs, f"Unknown layout: '{layout}'",
                       f"Block id={block.get('id','?')} in {block.get('_src','?')}. "
                       f"Known layouts: {', '.join(sorted(LAYOUTS))}.")
        return
    try:
        fn(prs, block)
    except Exception as exc:  # never let one bad block kill the build
        sys.stderr.write(f"[warn] failed to render '{layout}' "
                         f"(id={block.get('id','?')}, {block.get('_src','?')}): {exc}\n")
        _warning_slide(prs, f"Could not render '{layout}' slide",
                       f"id={block.get('id','?')} · {block.get('_src','?')} · {exc}")


def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    PAGE["n"] = 0

    blocks = load_content()

    if not blocks:
        # Empty / partial content dir: still produce a runnable deck.
        sys.stderr.write("[warn] no content blocks found — emitting placeholder deck\n")
        layout_cover(prs, {
            "title": "Future-State Scripting — A Strategy for Dis-Chem",
            "subtitle": "Executive strategy deck (content pending — engine scaffold)",
            "date": "June 2026",
        })
        _warning_slide(prs, "No content yet",
                       "deliverables/_shared/exec-deck/content/part*.md is empty or "
                       "missing. Content agents author slides to SPEC §3; this engine "
                       "renders them. Re-run once content exists.")
    else:
        # Ensure the first slide is a cover; if content didn't lead with one, add it.
        if blocks[0].get("_layout") != "cover":
            layout_cover(prs, {
                "title": "Future-State Scripting — A Strategy for Dis-Chem",
                "subtitle": "Executive strategy deck",
                "date": "June 2026",
            })
        for block in blocks:
            _render_block(prs, block)

    prs.save(OUT)
    return OUT, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    path, count = build()
    size = os.path.getsize(path)
    print(f"Wrote {path}")
    print(f"Size: {size:,} bytes ({size/1024:.0f} KB)")
    print(f"Slides: {count}")
    print("PDF step:  libreoffice --headless --convert-to pdf "
          f"{os.path.relpath(path)} --outdir deliverables/")
